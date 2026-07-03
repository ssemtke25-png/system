"""
탭7: 행사 계획서 기반 AI 문서 자동생성
리팩토링: 2026-07 / PPT 디자인 v2 + 문서4종 AI티 제거 업그레이드
"""

import streamlit as st
import google.generativeai as genai
import zipfile, io, re, json, requests
from pathlib import Path
from copy import copy as _copy
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 상수 ──────────────────────────────────────────────────────────────
GEMINI_MODEL  = "gemini-2.5-flash-lite"
FONT          = "맑은 고딕"

# PPT 색상 팔레트 (v2 디자인 시스템)
PALETTES = {
    "네이비 골드 (공식)": {
        "bg_dark":  (0x22, 0x33, 0x5A), "bg_dark2": (0x33, 0x47, 0x74),
        "bg_light": (0xFA, 0xFA, 0xFB), "card":     (0xFF, 0xFF, 0xFF),
        "card_tint":(0xF0, 0xF3, 0xF7), "panel":    (0x5C, 0x6B, 0x82),
        "accent":   (0x2E, 0x74, 0xB5), "accent2":  (0xC9, 0xA2, 0x27),
        "text_lt":  (0xFF, 0xFF, 0xFF), "text_dk":  (0x16, 0x1E, 0x2E),
        "sub_lt":   (0xD6, 0xE1, 0xF2), "sub_dk":   (0x5B, 0x66, 0x78),
        "num":      (0x9A, 0xA3, 0xB3),
    },
    "딥틸 트러스트": {
        "bg_dark":  (0x15, 0x3A, 0x40), "bg_dark2": (0x22, 0x53, 0x59),
        "bg_light": (0xF5, 0xFA, 0xF9), "card":     (0xFF, 0xFF, 0xFF),
        "card_tint":(0xE7, 0xF3, 0xF1), "panel":    (0x4A, 0x6E, 0x72),
        "accent":   (0x02, 0x89, 0x90), "accent2":  (0xE0, 0x8E, 0x45),
        "text_lt":  (0xEC, 0xFB, 0xF8), "text_dk":  (0x0E, 0x22, 0x22),
        "sub_lt":   (0xC7, 0xE3, 0xE1), "sub_dk":   (0x51, 0x66, 0x66),
        "num":      (0x8F, 0xAB, 0xAB),
    },
    "버건디 품격": {
        "bg_dark":  (0x3A, 0x1A, 0x22), "bg_dark2": (0x52, 0x28, 0x32),
        "bg_light": (0xFB, 0xF6, 0xF6), "card":     (0xFF, 0xFF, 0xFF),
        "card_tint":(0xF6, 0xE9, 0xEA), "panel":    (0x7A, 0x53, 0x58),
        "accent":   (0xB0, 0x2E, 0x3D), "accent2":  (0xC9, 0x9A, 0x4A),
        "text_lt":  (0xFD, 0xF0, 0xF1), "text_dk":  (0x2A, 0x14, 0x16),
        "sub_lt":   (0xE0, 0xB5, 0xB9), "sub_dk":   (0x6B, 0x50, 0x52),
        "num":      (0xB8, 0x95, 0x98),
    },
}

# AI티 제거 공통 지침 (문서4종 프롬프트에 공통 삽입)
ANTI_AI_GUIDE = """
[문장 다양성 - 반드시 지켜라]
- 문장 길이를 불균일하게 섞어라. 15자 안팎의 짧은 문장과 40자 이상의 긴 문장을 번갈아 써라.
  모든 문장이 비슷한 길이면 기계가 쓴 것처럼 보인다.
- 종결어미를 다양하게: "~하였다/~했다", "~할 예정이다", "~한다", 명사형 종결("~함", 개조식에서)을
  섞어서 써라. "~습니다"체로만 문장을 3개 이상 연속시키지 마라.
- 병렬 나열 구조("A하고 B하며 C한다")를 문서 전체에서 2회 이상 반복하지 마라.

[다음 상투어·클리셰 절대 금지]
"또한", "이러한 노력을 통해", "다양한 프로그램을 통해", "그 어느 때보다", "적극적으로 나서",
"실질적인 도움이 될 것으로 기대", "새로운 전기를 마련", "앞으로도 최선을 다하겠습니다",
"많은 관심과 참여 부탁드립니다", "뜻깊은 자리", "성공적으로 개최", "만전을 기하", "박차를 가하"
→ 위 표현이 쓰고 싶어지면, 대신 구체적 사실(숫자·일정·담당 조직명)로 문장을 채워라.

[구체성]
- "성공적으로 개최됩니다" 같은 뭉뚱그린 평가 대신, 실제 인원·건수·일정·장소 등
  주어진 정보의 숫자와 고유명사를 문장 속에 자연스럽게 녹여써라.
- 근거 없는 미사여구(예: "역사적인", "혁신적인")는 요약 정보에 실제로 그런 근거가
  있을 때만 사용하라.

[문단 연결]
- 문단을 "이를 통해", "그 결과" 같은 뻔한 연결어로 매번 맺지 마라.
  문단마다 다른 방식(질문형, 사실 제시형, 인과 없이 사실 나열 등)으로 넘어가라.
"""

# 보도자료 내장 예시 (참고파일 없을 때 폴백)
_PRESS_EXAMPLE = """
【보 도 자 료】
담당부서: 공간정보제도과  담당자: 홍길동  연락처: 044-000-0000

○○부, 지적재조사 담당자 역량강화 교육 실시
- 전국 시·군·구 담당 공무원 150명 대상, 최신 측량기술 집중 교육

국토교통부(장관 ○○○)는 10일 정부세종청사에서 전국 시·군·구 지적재조사 담당 공무원 150명을 대상으로 역량강화 교육을 실시했다.

이번 교육은 지적재조사사업의 현장 추진력을 높이기 위해 마련됐다. 드론 측량, 3D 공간정보 활용 등 최신 기술을 중심으로 진행됐으며, 우수사례 발표와 현장 실습도 병행됐다.

○○○ 국토교통부 지적재조사단장은 "이번 교육이 현장 담당자들의 실무 역량을 높이고 사업 추진에 실질적인 도움이 되길 바란다"고 말했다.

문의: 국토교통부 공간정보제도과 (044-000-0000)
"""

# ── AI 공통 ───────────────────────────────────────────────────────────
def get_model():
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    return genai.GenerativeModel(GEMINI_MODEL)

def ai(prompt: str) -> str:
    return get_model().generate_content(prompt).text.strip()

# ── HWPX 처리 ─────────────────────────────────────────────────────────
def hwpx_to_text(f) -> str:
    parts = []
    try:
        with zipfile.ZipFile(io.BytesIO(f.read()), "r") as zf:
            targets = sorted([n for n in zf.namelist()
                              if n.endswith(".xml") and ("Contents" in n or "content" in n.lower())])
            if not targets:
                targets = [n for n in zf.namelist() if n.endswith(".xml")]
            for name in targets:
                try:
                    raw = zf.read(name).decode("utf-8", errors="ignore")
                    clean = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", raw)).strip()
                    if len(clean) > 50:
                        parts.append(clean)
                except Exception:
                    pass
    except Exception as e:
        return f"[읽기 오류: {e}]"
    full = " ".join(parts)
    return full[:8000]

def summarize_hwpx(raw: str) -> str:
    return ai(f"""너는 20년 차 베테랑 공무원이다.
아래 행사 계획서 원문을 분석해 핵심 정보를 순수 JSON으로만 반환하라 (마크다운 없이).

원문: {raw}

JSON 형식:
{{"행사명":"","주최기관":"","주관기관":"","일시":"","장소":"","대상":"",
"참석예정인원":"","행사목적":"","주요프로그램":[""],
"예산":"","담당부서":"","담당자":"","기타특이사항":""}}""")

# ── 공통 UI ───────────────────────────────────────────────────────────
def show_textarea(sess_key: str, label: str, height: int = 420):
    if st.session_state.get(sess_key):
        st.caption("📋 텍스트 박스 클릭 → Ctrl+A → Ctrl+C → 한글/워드에 붙여넣기")
        st.text_area(label, value=st.session_state[sess_key],
                     height=height, key=f"ta_{sess_key}")

def gen_button(label: str, key: str, prompt: str, sess_key: str, height: int = 420):
    if st.button(label, key=key):
        with st.spinner("생성 중..."):
            st.session_state[sess_key] = ai(prompt)
    show_textarea(sess_key, label.replace("✍️ ", ""), height)

# ── 파일 텍스트 추출 (보도자료 참고용) ──────────────────────────────────
def extract_file_text(f) -> str:
    """txt / hwpx / pdf 에서 텍스트 추출 (최대 2,000자 — 보도자료 약 2장 분량)"""
    name = f.name.lower()
    try:
        if name.endswith(".txt"):
            return f.read().decode("utf-8", errors="ignore")[:2000]
        elif name.endswith((".hwpx", ".hwp")):
            return hwpx_to_text(f)[:2000]
        elif name.endswith(".pdf"):
            try:
                import fitz
                doc = fitz.open(stream=f.read(), filetype="pdf")
                return "\n".join(p.get_text() for p in doc)[:2000]
            except ImportError:
                return f.read().decode("utf-8", errors="ignore")[:2000]
    except Exception as e:
        return f"[추출 실패: {e}]"
    return ""

# ── 프롬프트 빌더 ─────────────────────────────────────────────────────
def _summary_str() -> str:
    return json.dumps(st.session_state.get("plan_summary_dict", {}),
                      ensure_ascii=False, indent=2)

def _prompt_report(s):
    return f"""너는 20년 차 베테랑 공무원이다. 신뢰감 있고 건조한 행정 공문서 문체로 작성하라.
감성적 수식어 배제. 개조식(○ 기호), 마크다운 없이 순수 텍스트.

{ANTI_AI_GUIDE}

[포함 항목] ○ 행사 개요 ○ 추진 목적 ○ 주요 내용 ○ 참석 대상·규모 ○ 소요 예산 ○ 기대 효과 ○ 향후 계획

행사 계획서 요약: {s}"""

def _prompt_director(s):
    return f"""너는 20년 차 베테랑 공무원이다. 국장급 인사말씀을 작성하라.
- 행사와 기관 정책 방향의 연계성 강조
- 내빈에게 정중하되 건조하고 신뢰감 있는 문체
- 서두 인삿말 포함 (따뜻한 어투 허용), 400~600자, 마크다운 없이

{ANTI_AI_GUIDE}

[구성] 1.서두 인삿말 2.행사의 행정적 의미·정책 연계 3.기관 비전 제시 4.마무리 당부

행사 계획서 요약: {s}"""

def _prompt_manager(s):
    return f"""너는 20년 차 베테랑 공무원이다. 과장급 인사말씀을 작성하라.
- 실무 총괄 과장 버전, 행사 준비 상황 언급 포함
- 참가자 실무적 당부·협조 요청
- 서두 인삿말 포함 (따뜻한 어투 허용), 300~500자, 마크다운 없이

{ANTI_AI_GUIDE}

[구성] 1.서두 2.준비 과정·노고 치하 3.목적·주요 내용 안내 4.당부·협조 요청 5.마무리

행사 계획서 요약: {s}"""

def _get_style_analysis(ref_texts: list) -> str:
    """참고자료 문체 특징만 추출 (1차 AI 호출 - 내용 오염 방지)"""
    examples = "\n\n---\n\n".join(
        f"[참고자료 {i+1}]\n{t}" for i, t in enumerate(ref_texts))
    return ai(
        "아래 보도자료들의 문체·형식 특징만 분석하라.\n"
        "내용(기관명·사업명·주제 등)은 절대 언급하지 말고 '쓰는 방식'만 정리하라.\n\n"
        "분석 항목:\n"
        "1. 문장 길이와 호흡\n"
        "2. 자주 쓰는 문장 종결 패턴\n"
        "3. 단락 구성 방식\n"
        "4. 제목/부제 형식\n"
        "5. 인용구 표현 방식\n"
        "6. 기타 문체 특징\n\n"
        f"참고자료:\n{examples}"
    )


def _prompt_press(s: str, style_analysis: str = "") -> str:
    """보도자료 생성 프롬프트 (style_analysis는 미리 추출된 문체 분석 결과)"""
    if style_analysis:
        style_section = f"[문체 분석 결과 - 이 스타일로 작성하라]\n{style_analysis}"
    else:
        style_section = (
            "[참고 문체]\n"
            "- '~했다', '~한다', '~이다' 기사체\n"
            "- 첫 문단은 육하원칙으로 간결하게\n"
            "- 단락은 3~5문장, 논리적 순서\n"
            "- 인용구: 직책+이름+\"...\"+이라고 말했다\n"
            "- 제목: 기관명+핵심행사+효과 구조"
        )
    return (
        "너는 20년 차 베테랑 공무원이자 언론홍보 전문가다.\n"
        "아래 [문체 스타일]로, 아래 [행사 계획서] 내용으로만 보도자료를 작성하라.\n\n"
        f"{style_section}\n\n"
        f"{ANTI_AI_GUIDE}\n\n"
        f"[행사 계획서 요약 - 오직 이 내용으로 보도자료를 작성할 것]\n{s}\n\n"
        "[필수 구조]\n"
        "1) 담당부서 / 담당자 / 연락처\n"
        "2) 제목: 위 행사의 핵심 한 문장\n"
        "3) 부제: 대시(-) 핵심 포인트 1~2개\n"
        "4) 본문: 육하원칙 → 추진 배경 → 목적 → 기대효과\n"
        "5) 기관장 인용구\n"
        "6) 문의처\n\n"
        "마크다운 없이 순수 텍스트로 작성하라."
    )

# ── 1. 문서 4종 ───────────────────────────────────────────────────────
def render_doc4():
    st.subheader("📄 문서 4종 자동생성")
    st.caption("생성 후 Ctrl+A → Ctrl+C → 한글/워드에 붙여넣기")
    s = _summary_str()

    with st.expander("📎 보도자료 참고파일 업로드 (선택 · 품질 향상)", expanded=False):
        st.caption("우리 기관 실제 보도자료를 올리면 문체·형식을 그대로 따라씁니다.")
        ref_files = st.file_uploader(
            "참고 보도자료 (txt / hwpx / pdf, 최대 3개)",
            type=["txt", "hwpx", "hwp", "pdf"],
            accept_multiple_files=True,
            key="press_ref_files",
        )
        if ref_files:
            cnt = min(len(ref_files), 3)
            st.success(f"✅ {cnt}개 반영됨 (3개 초과 시 앞 3개만 사용)")

    fkey = str([f.name for f in ref_files] if ref_files else [])
    if ref_files and st.session_state.get("_press_fkey") != fkey:
        with st.spinner("참고 보도자료 추출 중..."):
            st.session_state["press_ref_texts"] = [
                t for f in ref_files[:3]
                if (t := extract_file_text(f)) and len(t) > 50
            ]
        st.session_state["_press_fkey"] = fkey
    elif not ref_files:
        st.session_state["press_ref_texts"] = []

    ref_texts = st.session_state.get("press_ref_texts", [])

    col1, col2 = st.columns(2)
    with col1:
        gen_button("✍️ 요약보고서 생성", "btn_report",  _prompt_report(s),   "doc_report")
        st.markdown("---")
        gen_button("✍️ 과장인사말 생성", "btn_manager", _prompt_manager(s),  "doc_manager")
    with col2:
        gen_button("✍️ 국장인사말 생성", "btn_director", _prompt_director(s), "doc_director")
        st.markdown("---")
        st.markdown("**보도자료**")
        if st.button("✍️ 보도자료 생성", key="btn_press"):
            if ref_texts:
                with st.spinner("1단계: 참고자료 문체 분석 중..."):
                    style_analysis = _get_style_analysis(ref_texts)
                with st.spinner("2단계: 보도자료 작성 중..."):
                    st.session_state["doc_press"] = ai(_prompt_press(s, style_analysis))
            else:
                with st.spinner("보도자료 생성 중..."):
                    st.session_state["doc_press"] = ai(_prompt_press(s))
        show_textarea("doc_press", "보도자료", height=500)

# ── 2. 사회자 멘트 ────────────────────────────────────────────────────
MC_DEFAULT = ["개회선언","국민의례","내빈소개","기관장인사말","축사","주요프로그램 소개","폐회선언"]

def _prompt_mc(s, order, tone):
    return f"""너는 20년 차 베테랑 공무원이자 공공기관 전문 사회자다.
[톤]: {tone}  [순서]: {', '.join(order)}

[원칙]
- 각 순서는 ## 헤딩으로 구분
- 현장에서 바로 읽을 수 있는 완성된 문장, 각 100~200자
- 다음 순서로 자연스럽게 넘어가는 전환 문구 포함
- 내빈소개: 직책 → 성함 순

행사 계획서 요약: {s}"""

def render_mc():
    st.subheader("🎤 사회자 멘트 자동생성")
    s = _summary_str()
    with st.expander("⚙️ 순서 편집", expanded=False):
        edited = st.text_area("순서 (한 줄에 하나)", "\n".join(MC_DEFAULT),
                              height=200, key="mc_order_edit")
        order = [l.strip() for l in edited.split("\n") if l.strip()]
    st.info("순서: " + " → ".join(order))
    tone = st.selectbox("멘트 톤", ["격식체 (공식 행사)","친근체 (소규모/내부 행사)","방송체 (대규모/공개 행사)"])
    gen_button("✍️ 사회자 멘트 생성", "btn_mc", _prompt_mc(s, order, tone), "doc_mc", height=550)

# ── 3. 현수막 ─────────────────────────────────────────────────────────
def _prompt_banner(s, n, style):
    return f"""너는 20년 차 베테랑 공무원이자 공공기관 홍보 전문가다.
현수막 문구 시안 {n}개를 작성하라.
[스타일]: {style}
[원칙] 15~25자 이내 / 번호. 문구 형식 / 각 문구 아래 사용 설명 한 줄 / 모호한 표현 배제
행사 계획서 요약: {s}"""

def render_banner():
    st.subheader("🪧 현수막 문안")
    summary = st.session_state.get("plan_summary_dict", {})

    if summary:
        fields = [
            ("행  사  명", summary.get("행사명", "")),
            ("기      간", summary.get("일시", "")),
            ("장      소", summary.get("장소", "")),
            ("주      최", summary.get("주최기관", "")),
            ("주      관", summary.get("주관기관", "")),
            ("대      상", summary.get("대상", "")),
            ("담 당 부 서", summary.get("담당부서", "")),
        ]
        lines = [f"{k} : {v}" for k, v in fields if v and str(v).strip()]
        banner_text = "■ 현수막 문안\n\n" + "\n".join(lines)

        st.caption("📋 계획서에서 자동 추출 — 복사하거나 .txt로 다운로드하세요.")
        st.text_area("현수막 문안", value=banner_text, height=260, key="ta_banner_auto")
        st.download_button(
            "⬇️ .txt 다운로드",
            data=banner_text.encode("utf-8"),
            file_name=f"현수막문안_{summary.get('행사명','행사')}.txt",
            mime="text/plain",
        )
    else:
        st.info("HWPX 계획서를 업로드하면 현수막 문안이 자동으로 추출됩니다.")

    st.markdown("---")
    st.markdown("**💡 홍보 슬로건 생성 (선택)**")
    st.caption("행사명 외에 현수막에 들어갈 임팩트 있는 슬로건이 필요할 때 사용하세요.")
    s = _summary_str()
    c1, c2 = st.columns(2)
    n     = c1.slider("시안 수", 3, 7, 5)
    style = c2.selectbox("스타일", ["공식/격식형","친근/따뜻형","역동/강조형","혼합 (다양하게)"])
    gen_button("✍️ 슬로건 생성", "btn_banner", _prompt_banner(s, n, style), "doc_banner", height=280)

# ── 4. 결과보고서 ─────────────────────────────────────────────────────
def _prompt_result(s, attendance, satisfaction, note):
    return f"""너는 20년 차 베테랑 공무원이다. 행사 결과보고서 초안을 작성하라.
감성 배제, 수치·사실 위주, 개조식(○), 마크다운 없이.

[결과 정보] 참석인원:{attendance or '미기재'} / 만족도:{satisfaction or '미기재'} / 특이사항:{note or '없음'}

[형식]
1.행사 개요 (○ 행사명/일시/장소/주최·주관/참석인원)
2.추진 결과 (○ 참석현황 계획대비 / 프로그램별 결과 / 만족도)
3.예산 집행 (○ 편성/집행/집행률)
4.주요 성과 및 시사점 (○ 성과 / 미흡사항·개선방향)
5.향후 계획 (○ 후속조치)

행사 계획서 요약: {s}"""

def render_result():
    st.subheader("📋 결과보고서 초안")
    s = _summary_str()
    c1, c2 = st.columns(2)
    attendance   = c1.text_input("실제 참석 인원",    placeholder="예: 150명")
    satisfaction = c1.text_input("만족도 조사 결과",  placeholder="예: 4.2/5.0")
    note         = c2.text_area("주요 성과·특이사항", placeholder="예) 전년 대비 참석률 20% 증가", height=105)
    gen_button("✍️ 결과보고서 초안 생성", "btn_result",
               _prompt_result(s, attendance, satisfaction, note), "doc_result", height=550)

# ── 6. 명찰 ──────────────────────────────────────────────────────────
def render_namecard():
    st.subheader("🪪 명찰 자동생성")
    summary = st.session_state.get("plan_summary_dict", {})
    event_name = summary.get("행사명","행사") if summary else "행사"
    st.info("엑셀 파일에 **이름**, **소속** 컬럼이 있어야 합니다.")
    c1,c2=st.columns(2)
    excel_file=c1.file_uploader("명단 엑셀 업로드",type=["xlsx","xls"],key="nc_excel")
    schedule=c2.text_area("뒷면 일정표",placeholder="09:00 등록\n10:00 개회식",height=150)
    ev=st.text_input("행사명 (비워두면 자동)",placeholder=event_name) or event_name

    if excel_file and st.button("🪪 명찰 생성"):
        try:
            from openpyxl import load_workbook
            wb=load_workbook(io.BytesIO(excel_file.read()),data_only=True); ws=wb.active
            hdrs=[str(c.value).strip() if c.value else "" for c in ws[1]]
            ni=next((i for i,h in enumerate(hdrs) if any(k in h for k in ["이름","성명","name","Name"])),None)
            oi=next((i for i,h in enumerate(hdrs) if any(k in h for k in ["소속","기관","org","Org","dept"])),None)
            if ni is None: st.error(f"이름 컬럼 없음. 현재: {hdrs}"); return
            persons=[{"name":str(r[ni]).strip(),"org":str(r[oi]).strip() if oi and r[oi] else ""}
                     for r in ws.iter_rows(min_row=2,values_only=True)
                     if r[ni] and str(r[ni]).lower()!="none"]
            if not persons: st.error("유효한 데이터 없음"); return
            st.info(f"총 {len(persons)}명 처리 중...")
            data=_build_namecard(persons,ev,schedule)
            if data:
                st.session_state["nc_bytes"]=data
                st.session_state["nc_count"]=len(persons)
        except Exception as e:
            st.error(f"오류: {e}")
            import traceback; st.text(traceback.format_exc())

    if st.session_state.get("nc_bytes"):
        st.success(f"✅ {st.session_state['nc_count']}명 생성 완료!")
        st.download_button(f"⬇️ 명찰 다운로드",data=st.session_state["nc_bytes"],
                           file_name="명찰.docx",
                           mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

def _build_namecard(persons, event_name, schedule_text):
    try:
        from docx import Document
        from docx.shared import Pt,Cm,RGBColor as DocxRGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc=Document()
        for sec in doc.sections:
            sec.page_width=Cm(21); sec.page_height=Cm(29.7)
            sec.top_margin=sec.bottom_margin=Cm(1.5)
            sec.left_margin=sec.right_margin=Cm(2.0)

        def dashed(d):
            hr=d.add_paragraph(); hr.paragraph_format.space_before=hr.paragraph_format.space_after=Pt(0)
            pPr=hr._p.get_or_add_pPr(); pBdr=OxmlElement("w:pBdr")
            b=OxmlElement("w:bottom")
            b.set(qn("w:val"),"dashed"); b.set(qn("w:sz"),"6")
            b.set(qn("w:space"),"1"); b.set(qn("w:color"),"888888")
            pBdr.append(b); pPr.append(pBdr)

        def para(d,text,size,bold=False,color=(0x1A,0x1A,0x2E),before=4,after=4,align=WD_ALIGN_PARAGRAPH.CENTER):
            p=d.add_paragraph(); p.alignment=align
            p.paragraph_format.space_before=Pt(before); p.paragraph_format.space_after=Pt(after)
            r=p.add_run(text); r.font.name=FONT; r.font.size=Pt(size)
            r.font.bold=bold; r.font.color.rgb=DocxRGBColor(*color)

        def front(name,org,first):
            if not first: dashed(doc)
            para(doc,event_name,15,True,(0x1F,0x38,0x64),28,8)
            if org: para(doc,org,17,False,(0x2E,0x74,0xB5),4,4)
            para(doc,name,38,True,(0x1A,0x1A,0x2E),8,28)

        def back(first):
            if not first: dashed(doc)
            para(doc,"행사 일정",15,True,(0x1F,0x38,0x64),28,10)
            for line in (schedule_text or "일정 미정").strip().split("\n"):
                para(doc,line.strip(),13,False,(0x1A,0x1A,0x2E),3,3)
            para(doc,"",13,False,(0x1A,0x1A,0x2E),4,28)

        for i in range(0,len(persons),2):
            batch=persons[i:i+2]
            front(batch[0]["name"],batch[0]["org"],True)
            if len(batch)>1: front(batch[1]["name"],batch[1]["org"],False)
            doc.add_page_break()
            back(True)
            if len(batch)>1: back(False)
            if i+2<len(persons): doc.add_page_break()

        buf=io.BytesIO(); doc.save(buf); return buf.getvalue()
    except ImportError:
        st.error("pip install python-docx"); return None
    except Exception as e:
        st.error(f"명찰 오류: {e}")
        import traceback; st.text(traceback.format_exc()); return None

# ── 7. 행사장 약도 ────────────────────────────────────────────────────
def render_map():
    st.subheader("🗺️ 행사장 약도")
    summary=st.session_state.get("plan_summary_dict",{})
    default=summary.get("장소","") if summary else ""
    place=st.text_input("장소명 또는 주소",value=default,
                        placeholder="예: 부산 벡스코  또는  부산광역시 해운대구 APEC로 55")
    c1,c2,c3=st.columns(3)
    zoom=c1.slider("확대 수준",1,14,4,help="숫자가 클수록 확대")
    mw=c2.number_input("가로(px)",200,1200,800)
    mh=c3.number_input("세로(px)",200,900,600)

    if st.button("🗺️ 약도 생성") and place:
        try:
            kkey=st.secrets["KAKAO_API_KEY"]
            hdrs={"Authorization":f"KakaoAK {kkey}"}

            def kw(q):
                r=requests.get("https://dapi.kakao.com/v2/local/search/keyword.json",
                               headers=hdrs,params={"query":q})
                d=r.json().get("documents",[])
                return d[0] if d else None

            def addr(q):
                r=requests.get("https://dapi.kakao.com/v2/local/search/address.json",
                               headers=hdrs,params={"query":q})
                d=r.json().get("documents",[])
                return d[0] if d else None

            pm=re.search(r"\(([^)]+)\)",place)
            paddr=pm.group(1).strip() if pm else None
            clean=re.sub(r"\s*\([^)]*\)","",place).strip()

            doc0=None
            if paddr:
                st.info(f"🔍 주소 감지: **{paddr}**")
                a=addr(paddr)
                if a:
                    an=((a.get("road_address") or {}).get("address_name")
                        or (a.get("address") or {}).get("address_name",paddr))
                    doc0={"x":a["x"],"y":a["y"],"place_name":clean,"road_address_name":an}

            if not doc0: doc0=kw(clean)
            if not doc0: doc0=kw(place)
            if not doc0:
                with st.spinner("AI가 주소 보정 중..."):
                    fixed=ai(f"다음 장소명의 정확한 도로명 주소를 한 줄로만 출력하세요.\n장소명: {clean}")
                st.info(f"🔍 AI 보정: **{fixed}**")
                doc0=kw(fixed)
                if not doc0:
                    a=addr(fixed)
                    if a:
                        an=((a.get("road_address") or {}).get("address_name")
                            or (a.get("address") or {}).get("address_name",fixed))
                        doc0={"x":a["x"],"y":a["y"],"place_name":clean,"road_address_name":an}

            if not doc0 or not doc0.get("x"):
                st.error("위치를 찾을 수 없습니다. 도로명 주소를 직접 입력해보세요.")
                return

            lng,lat=float(doc0["x"]),float(doc0["y"])
            fn=doc0.get("place_name",clean)
            fa=doc0.get("road_address_name") or doc0.get("address_name","")
            st.success(f"📍 찾은 장소: **{fn}** ({fa})")

            mr=requests.get("https://dapi.kakao.com/v2/maps/staticmap",headers=hdrs,
                            params={"center":f"{lng},{lat}","level":zoom,
                                    "w":int(mw),"h":int(mh),"markers":f"color:red|{lng},{lat}"})
            ct=mr.headers.get("Content-Type","")
            if mr.status_code==200 and "image" in ct:
                st.session_state.update({"map_img":mr.content,"map_name":fn})
            else:
                ourl=(f"https://staticmap.openstreetmap.de/staticmap.php"
                      f"?center={lat},{lng}&zoom={zoom+2}&size={int(mw)}x{int(mh)}"
                      f"&markers={lat},{lng},red-pushpin")
                or_=requests.get(ourl,timeout=10)
                if or_.status_code==200 and "image" in or_.headers.get("Content-Type",""):
                    st.session_state.update({"map_img":or_.content,"map_name":fn})
                    st.caption("※ OpenStreetMap 기반")
                else:
                    from urllib.parse import quote
                    st.warning("지도 이미지 생성 실패.")
                    st.markdown(f"[🗺️ 카카오맵에서 보기](https://map.kakao.com/?q={quote(fa or clean)})")
        except KeyError:
            st.error("KAKAO_API_KEY가 secrets에 없습니다.")
        except Exception as e:
            st.error(f"오류: {e}")
            import traceback; st.text(traceback.format_exc())

    if st.session_state.get("map_img"):
        fn=st.session_state.get("map_name","행사장")
        st.image(st.session_state["map_img"],caption=f"📍 {fn}",use_container_width=True)
        st.download_button("⬇️ 약도 이미지 다운로드",data=st.session_state["map_img"],
                           file_name=f"행사장약도_{fn}.png",mime="image/png")

# ── 5. PPT (python-pptx) v2 디자인 ────────────────────────────────────
def _prompt_ppt(s, n):
    return (
        "너는 20년 차 베테랑 공무원이자 공공기관 발표자료 전문가다.\n"
        "감성 배제, 수치·사실 위주, 담당자가 바로 발표할 수 있는 수준으로 작성하라.\n\n"
        "[레이아웃 7종 - 반드시 아래 필드를 빠짐없이 채워라]\n"
        "title   : 표지       → title(행사명), subtitle(일시|장소|주최)\n"
        "closing : 마무리     → title(감사합니다), subtitle(담당부서)\n"
        "section : 챕터구분   → title만\n"
        "content : 일반슬라이드 → title, body(2~3문장 본문), bullets(3~5개 필수)\n"
        "two_column: 좌우비교 → title, body, left_title, left_bullets(3개↑), right_title, right_bullets(3개↑)\n"
        "highlight : 숫자강조 → title, body, stat_number(숫자), stat_label(설명), bullets(3개↑)\n"
        "table   : 표         → title, body, headers(컬럼명[]), rows(데이터[][], 3행↑)\n\n"
        "[필수 슬라이드 순서]\n"
        "표지(title) → 목차(content, '01. 02.' 형식 bullets 4~6개) → "
        "행사개요(highlight, 참석인원 stat_number) → 추진배경(two_column, 현황vs목표) → "
        "주요프로그램(table, 시간표 3행↑) → 세부내용(content) → 기대효과(content) → 클로징(closing)\n\n"
        "⚠️ bullets/left_bullets/right_bullets/rows 가 비어있는 슬라이드 절대 금지\n"
        "⚠️ body 필드 없는 슬라이드 절대 금지 (title/closing/section 제외)\n\n"
        "순수 JSON 배열만 반환 (마크다운 없이).\n\n"
        f"슬라이드 수: {n}개\n행사 계획서 요약:\n{s}"
    )


def _is_empty(si: dict) -> bool:
    lay = si.get("layout", "content")
    if lay in ("title", "closing", "section"):
        return False
    if lay == "table":
        return not si.get("rows") or len(si.get("rows", [])) == 0
    if lay == "two_column":
        return not si.get("left_bullets") and not si.get("right_bullets")
    if lay == "highlight":
        return not si.get("bullets") and not si.get("stat_number")
    return not si.get("bullets") or len(si.get("bullets", [])) == 0


def _validate_and_fix(slides, summary_str):
    needs_fix = [i for i, si in enumerate(slides) if _is_empty(si)]
    if not needs_fix:
        return slides
    targets = [slides[i] for i in needs_fix]
    fix_prompt = (
        "아래 슬라이드들의 내용이 비어있다. 각 슬라이드 layout에 맞게 내용을 채워라.\n"
        "- content → bullets 3~5개 필수\n"
        "- table → rows 3행 이상 필수\n"
        "- two_column → left_bullets, right_bullets 각 3개 이상\n"
        "- highlight → stat_number, stat_label, bullets 3개 이상\n"
        "JSON 배열로만 반환 (마크다운 없이, 슬라이드 수 유지).\n\n"
        f"행사 계획서 요약:\n{summary_str}\n\n"
        f"보완할 슬라이드:\n{json.dumps(targets, ensure_ascii=False)}"
    )
    try:
        raw = ai(fix_prompt)
        fixed = json.loads(re.sub(r"```json|```", "", raw).strip())
        for idx, fi in zip(needs_fix, fixed):
            slides[idx] = fi
    except Exception:
        pass
    return slides


def render_ppt():
    st.subheader("📊 PPT 자동생성")
    st.caption("행사 계획서 기반으로 현장에서 바로 발표 가능한 PPT를 생성합니다.")
    s = _summary_str()
    c1, c2 = st.columns([1, 2])
    n     = c1.slider("슬라이드 수", 8, 20, 12)
    theme = c2.selectbox("색상 테마", list(PALETTES.keys()))

    with st.expander("🖼️ 표지·클로징용 이미지 업로드 (선택)", expanded=False):
        st.caption("행사 포스터나 대표 사진을 올리면 표지/클로징 우측 절반에 실제 이미지가 들어갑니다. 없으면 자동 생성 그래픽이 대신 들어갑니다.")
        cover_img = st.file_uploader("이미지 (jpg/png)", type=["jpg", "jpeg", "png"], key="ppt_cover_img")
    cover_bytes = cover_img.read() if cover_img else None

    if st.button("🖥️ PPT 생성", type="primary"):
        with st.spinner("AI가 슬라이드 구성 중... (10~20초)"):
            raw = ai(_prompt_ppt(s, n))
        try:
            slides = json.loads(re.sub(r"```json|```", "", raw).strip())
        except Exception as e:
            st.error(f"파싱 오류: {e}")
            with st.expander("AI 원본 확인"):
                st.text(raw[:1000])
            return

        empty_cnt = sum(1 for si in slides if _is_empty(si))
        if empty_cnt:
            with st.spinner(f"⚠️ 빈 슬라이드 {empty_cnt}개 감지 → 자동 보완 중..."):
                slides = _validate_and_fix(slides, s)
            still_empty = sum(1 for si in slides if _is_empty(si))
            if still_empty:
                st.warning(f"⚠️ {still_empty}개 슬라이드는 내용을 직접 채워주세요.")

        with st.spinner("PPT 파일 생성 중..."):
            data = _build_pptx(slides, st.session_state.get("plan_summary_dict", {}), theme, cover_bytes)
        if data:
            name = st.session_state.get("plan_summary_dict", {}).get("행사명", "행사")
            st.session_state.update({"ppt_bytes": data,
                                     "ppt_name": f"{name}_발표자료.pptx",
                                     "ppt_count": len(slides)})

    if st.session_state.get("ppt_bytes"):
        st.success(f"✅ {st.session_state['ppt_count']}개 슬라이드 생성 완료!")
        st.download_button(
            "⬇️ PPT 다운로드 (.pptx)",
            data=st.session_state["ppt_bytes"],
            file_name=st.session_state["ppt_name"],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )


def _build_pptx(slides_data, summary, theme="네이비 골드 (공식)", cover_image_bytes=None):
    """
    v2 디자인: 안티패턴(좌측 스트라이프/제목밑줄/카드 단측테두리) 제거,
    원형 번호 뱃지 통일 모티프, 다크(표지·구분·클로징)/라이트(본문) 샌드위치 구조,
    목차는 Contents 분할형 전용 레이아웃, 표지·클로징은 이미지 업로드 시 half-bleed 반영.
    """
    from pptx import Presentation

    pal = PALETTES.get(theme, PALETTES["네이비 골드 (공식)"])
    def C(k): return RGBColor(*pal[k])
    def W():  return RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W_, H_ = prs.slide_width, prs.slide_height
    BL = prs.slide_layouts[6]

    def RECT(sl,x,y,w,h,k):
        s=sl.shapes.add_shape(1,x,y,w,h)
        s.fill.solid(); s.fill.fore_color.rgb=C(k)
        s.line.fill.background(); s.shadow.inherit=False
        return s
    def CARD(sl,x,y,w,h,k="card_tint"):
        s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
        try: s.adjustments[0]=0.045
        except Exception: pass
        s.fill.solid(); s.fill.fore_color.rgb=C(k)
        s.line.fill.background(); s.shadow.inherit=False
        return s
    def TXT(sl,text,x,y,w,h,sz,bold=False,k="text_dk",al=PP_ALIGN.LEFT,it=False):
        tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=al
        r=p.add_run(); r.text=str(text)
        r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=it
        r.font.name=FONT; r.font.color.rgb=C(k)
        return tb
    def CIRCLE(sl,cx,cy,d,label,fill_k="accent",text_k=None):
        s=sl.shapes.add_shape(MSO_SHAPE.OVAL,cx,cy,d,d)
        s.fill.solid(); s.fill.fore_color.rgb=C(fill_k)
        s.line.fill.background(); s.shadow.inherit=False
        tf=s.text_frame; tf.word_wrap=False
        tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(label)
        r.font.size=Pt(max(10,int(Emu(d).inches*40)))
        r.font.bold=True; r.font.name=FONT
        r.font.color.rgb=W() if not text_k else C(text_k)
        return s
    def DOTS(sl,items,x,y,w,h,sz=15,dot_k="accent",text_k="text_dk",gap=10):
        if not items: return
        tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
        for i,b in enumerate(items):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.space_before=Pt(gap); p.space_after=Pt(2)
            r1=p.add_run(); r1.text="●  "
            r1.font.size=Pt(sz-3); r1.font.name=FONT; r1.font.color.rgb=C(dot_k); r1.font.bold=True
            r2=p.add_run(); r2.text=str(b)
            r2.font.size=Pt(sz); r2.font.name=FONT; r2.font.color.rgb=C(text_k)
    def PAGENUM(sl,n,k="sub_dk"):
        tb=sl.shapes.add_textbox(W_-Inches(0.8),H_-Inches(0.45),Inches(0.6),Inches(0.3))
        p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=str(n)
        r.font.size=Pt(10); r.font.name=FONT; r.font.color.rgb=C(k)
    def CIRCLE_PLAIN(sl,x,y,d,k,blend_bg=None):
        s=sl.shapes.add_shape(MSO_SHAPE.OVAL,x,y,d,d)
        if blend_bg:
            base=pal[k]; mixed=tuple(int(base[i]*0.35+blend_bg[i]*0.65) for i in range(3))
            s.fill.solid(); s.fill.fore_color.rgb=RGBColor(*mixed)
        else:
            s.fill.solid(); s.fill.fore_color.rgb=C(k)
        s.line.fill.background(); s.shadow.inherit=False
        return s
    def DECOR_PANEL(sl,px,py,pw,ph,image_bytes=None):
        if image_bytes:
            sl.shapes.add_picture(io.BytesIO(image_bytes),px,py,pw,ph); return
        RECT(sl,px,py,pw,ph,"bg_dark2")
        cx=px+int(pw*0.62); cy=py+int(ph*0.42)
        CIRCLE_PLAIN(sl,cx-Inches(2.6),cy-Inches(2.6),Inches(5.2),"accent2",pal["bg_dark2"])
        CIRCLE_PLAIN(sl,cx-Inches(1.6),cy-Inches(1.6),Inches(3.2),"accent",pal["bg_dark2"])
        CIRCLE_PLAIN(sl,cx-Inches(0.5),cy-Inches(0.5),Inches(1.0),"accent2")
    def slice_bullets(buls):
        return [str(b).lstrip("0123456789. ").strip() for b in buls]

    for idx, si in enumerate(slides_data):
        sl = prs.slides.add_slide(BL)
        lay  = str(si.get("layout", "content")).strip().lower()
        # AI가 layout 문자열을 정확히 안 맞춰 보내는 경우 대비 - 위치로 강제 보정
        if idx == 0:
            lay = "title"
        elif idx == len(slides_data) - 1:
            lay = "closing"
        ttl  = si.get("title", "")
        sub  = si.get("subtitle", "")
        body = si.get("body", "")
        buls = si.get("bullets", []) or []
        hdrs = si.get("headers", []) or []
        rows = si.get("rows", []) or []
        sn   = idx + 1
        is_toc = lay == "content" and any(str(b).startswith(("01","02","1.","2.")) for b in buls)

        if lay in ("title", "closing", "section"):
            RECT(sl,0,0,W_,H_,"bg_dark")
            panel_w = int(W_*0.42)
            DECOR_PANEL(sl, W_-panel_w, 0, panel_w, H_,
                        cover_image_bytes if lay in ("title","closing") else None)
            if lay == "title":
                org=(summary or {}).get("주최기관","")
                if org: TXT(sl,org,Inches(0.6),Inches(0.6),Inches(6.5),Inches(0.5),14,True,k="accent2")
                TXT(sl,ttl,Inches(0.6),Inches(2.3),Inches(7.0),Inches(2.6),38,True,k="text_lt")
                if sub: TXT(sl,sub,Inches(0.6),Inches(5.1),Inches(7.0),Inches(1.0),14,k="sub_lt")
            elif lay == "closing":
                TXT(sl,ttl,Inches(0.6),Inches(2.6),Inches(7.0),Inches(1.6),42,True,k="text_lt")
                if sub: TXT(sl,sub,Inches(0.6),Inches(3.9),Inches(7.0),Inches(0.7),16,k="sub_lt")
                dept=" ".join(filter(None,[(summary or {}).get("담당부서",""),(summary or {}).get("담당자","")])).strip()
                if dept: TXT(sl,dept,Inches(0.6),Inches(6.7),Inches(7.0),Inches(0.4),11,k="sub_lt",it=True)
            elif lay == "section":
                CIRCLE(sl,Inches(0.6),Inches(2.6),Inches(1.1),f"{sn}","accent2")
                TXT(sl,ttl,Inches(0.6),Inches(4.0),Inches(7.5),Inches(1.6),32,True,k="text_lt")

        elif is_toc:
            RECT(sl,0,0,W_,H_,"bg_light")
            left_w = int(W_*0.38)
            RECT(sl,0,0,left_w,H_,"panel")
            TXT(sl,"Contents",0,int(H_*0.44),left_w,Inches(0.9),36,True,k="text_lt",al=PP_ALIGN.CENTER)
            line_y = int(H_*0.44) + Inches(0.8)
            RECT(sl, int(left_w*0.30), line_y, int(left_w*0.40), Pt(1.2), "bg_light")

            items = slice_bullets(buls)[:6]
            right_x = left_w + Inches(0.5); right_w = W_ - left_w - Inches(0.9)
            row_h = min(Inches(0.98), (H_ - Inches(1.2)) // max(len(items),1))
            start_y = (H_ - row_h*len(items)) // 2
            for i, title in enumerate(items):
                ry = start_y + row_h*i
                TXT(sl,f"{i+1}",right_x,ry,Inches(0.9),Inches(0.8),36,True,k="num")
                TXT(sl,title,right_x+Inches(1.0),ry+Inches(0.14),right_w-Inches(1.0),Inches(0.55),19,True,k="text_dk")
                if i < len(items)-1:
                    RECT(sl,right_x,ry+row_h-Inches(0.06),right_w,Pt(1),"card_tint")
            PAGENUM(sl,sn)

        else:
            RECT(sl,0,0,W_,H_,"bg_light")
            CIRCLE(sl,Inches(0.5),Inches(0.45),Inches(0.62),f"{sn}","accent")
            TXT(sl,ttl,Inches(1.35),Inches(0.42),Inches(10.5),Inches(0.6),26,True,k="text_dk")
            if body:
                TXT(sl,body,Inches(1.35),Inches(1.05),Inches(11.2),Inches(0.6),12.5,k="sub_dk")
                content_top = Inches(1.85)
            else:
                content_top = Inches(1.35)

            if lay == "highlight":
                stat_n=si.get("stat_number",""); stat_l=si.get("stat_label","")
                card_w=Inches(4.3); card_h=H_-content_top-Inches(0.5)
                CARD(sl,Inches(0.5),content_top,card_w,card_h)
                TXT(sl,stat_n,Inches(0.5),content_top+Inches(0.55),card_w,Inches(1.5),52,True,k="accent",al=PP_ALIGN.CENTER)
                TXT(sl,stat_l,Inches(0.5),content_top+Inches(2.0),card_w,Inches(0.5),14,k="sub_dk",al=PP_ALIGN.CENTER)
                DOTS(sl,buls,Inches(5.2),content_top+Inches(0.1),Inches(7.6),card_h,14)
                PAGENUM(sl,sn)

            elif lay == "two_column":
                lt=si.get("left_title","현황"); lb=si.get("left_bullets",[]) or []
                rt=si.get("right_title","목표"); rb=si.get("right_bullets",[]) or []
                cw=Inches(5.85); ch=H_-content_top-Inches(0.5)
                lx=Inches(0.5); rx=lx+cw+Inches(0.3)
                for x,ct,cb,bk in [(lx,lt,lb,"accent"),(rx,rt,rb,"accent2")]:
                    CARD(sl,x,content_top,cw,ch)
                    CIRCLE(sl,x+Inches(0.25),content_top+Inches(0.25),Inches(0.42),"",bk)
                    TXT(sl,ct,x+Inches(0.85),content_top+Inches(0.22),cw-Inches(1.1),Inches(0.5),16,True,k="text_dk")
                    DOTS(sl,cb,x+Inches(0.35),content_top+Inches(0.95),cw-Inches(0.6),ch-Inches(1.1),13,dot_k=bk)
                PAGENUM(sl,sn)

            elif lay=="table" and hdrs and rows:
                cc=len(hdrs); rc=len(rows)+1
                tbl=sl.shapes.add_table(rc,cc,Inches(0.5),content_top,Inches(12.3),H_-content_top-Inches(0.5)).table
                cw_=Inches(12.3)//cc
                for c in range(cc): tbl.columns[c].width=cw_
                for c,h in enumerate(hdrs):
                    cell=tbl.cell(0,c); cell.text=str(h)
                    cell.fill.solid(); cell.fill.fore_color.rgb=C("accent")
                    p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
                    r=p.runs[0] if p.runs else p.add_run()
                    r.font.bold=True; r.font.size=Pt(12); r.font.color.rgb=W(); r.font.name=FONT
                for ri,row in enumerate(rows):
                    rk="card" if ri%2==0 else "card_tint"
                    for c,v in enumerate(row[:cc]):
                        cell=tbl.cell(ri+1,c); cell.text=str(v)
                        cell.fill.solid(); cell.fill.fore_color.rgb=C(rk)
                        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
                        r=p.runs[0] if p.runs else p.add_run()
                        r.font.size=Pt(11); r.font.color.rgb=C("text_dk"); r.font.name=FONT
                PAGENUM(sl,sn)

            else:
                DOTS(sl,buls,Inches(0.6),content_top,Inches(11.8),H_-content_top-Inches(0.5),15)
                PAGENUM(sl,sn)

    buf=io.BytesIO(); prs.save(buf); return buf.getvalue()


# ── 메인 ──────────────────────────────────────────────────────────────
def render_tab7():
    st.title("🤖 AI 문서 자동생성")
    st.caption("HWPX 계획서를 업로드하면 각종 문서를 자동으로 생성합니다.")
    st.markdown("---")
    st.subheader("📁 행사 계획서 업로드")

    hwpx=st.file_uploader("HWPX 파일 업로드 (.hwpx / .hwp)",type=["hwpx","hwp"])
    if hwpx:
        fkey=f"{hwpx.name}_{hwpx.size}"
        if (st.session_state.get("hwpx_fkey")!=fkey or "plan_summary_raw" not in st.session_state):
            st.session_state["hwpx_fkey"]=fkey
            with st.spinner("텍스트 추출 및 요약 중..."):
                raw=hwpx_to_text(hwpx)
                if len(raw)<100:
                    st.error("텍스트를 충분히 추출하지 못했습니다."); return
                summary_raw=summarize_hwpx(raw)
                st.session_state["plan_summary_raw"]=summary_raw
                try:
                    st.session_state["plan_summary_dict"]=json.loads(
                        re.sub(r"```json|```","",summary_raw).strip())
                except Exception:
                    st.session_state["plan_summary_dict"]={}

        summary=st.session_state.get("plan_summary_dict",{})
        with st.expander("📋 계획서 요약 (모든 기능에 재사용됨)",expanded=True):
            if summary:
                c1,c2=st.columns(2); keys=list(summary.keys()); half=len(keys)//2
                for i,k in enumerate(keys):
                    v=summary[k]
                    if isinstance(v,list): v=", ".join(str(x) for x in v)
                    (c1 if i<half else c2).markdown(f"**{k}:** {v}")
            else:
                st.text(st.session_state.get("plan_summary_raw",""))
    else:
        st.info("HWPX 업로드 없이도 수동으로 사용 가능합니다.")
        if "plan_summary_dict" not in st.session_state:
            st.session_state["plan_summary_dict"]={}

    st.markdown("---")
    tabs=st.tabs(["📄 문서 4종","🎤 사회자 멘트","🪧 현수막","📋 결과보고서","📊 PPT","🪪 명찰","🗺️ 행사장 약도"])
    with tabs[0]: render_doc4()
    with tabs[1]: render_mc()
    with tabs[2]: render_banner()
    with tabs[3]: render_result()
    with tabs[4]: render_ppt()
    with tabs[5]: render_namecard()
    with tabs[6]: render_map()

render = render_tab7  # main.py 호환

if __name__ == "__main__":
    st.set_page_config(page_title="AI 문서 자동생성", page_icon="🤖", layout="wide")
    render_tab7()
