import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FONT = "맑은 고딕"

PALETTES = {
    "네이비 골드 (공식)": {
        "bg_dark":  (0x22, 0x33, 0x5A), "bg_dark2": (0x33, 0x47, 0x74),
        "bg_light": (0xFA, 0xFA, 0xFB), "card":     (0xFF, 0xFF, 0xFF),
        "card_tint":(0xF0, 0xF3, 0xF7), "panel":    (0x5C, 0x6B, 0x82),
        "accent":   (0x2E, 0x74, 0xB5), "accent2":  (0xC9, 0xA2, 0x27),
        "text_lt":  (0xFF, 0xFF, 0xFF), "text_dk":  (0x16, 0x1E, 0x2E),
        "sub_lt":   (0xD6, 0xE1, 0xF2), "sub_dk":   (0x5B, 0x66, 0x78),
        "num":      (0x9A, 0xA3, 0xB3),
    },
}


def _build_pptx(slides_data, summary, theme="네이비 골드 (공식)", cover_image_bytes=None):
    pal = PALETTES.get(theme, PALETTES["네이비 골드 (공식)"])
    def C(k): return RGBColor(*pal[k])
    def W(): return RGBColor(0xFF,0xFF,0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    W_, H_ = prs.slide_width, prs.slide_height
    BL = prs.slide_layouts[6]

    def RECT(sl,x,y,w,h,k):
        s=sl.shapes.add_shape(1,x,y,w,h)
        s.fill.solid(); s.fill.fore_color.rgb=C(k)
        s.line.fill.background(); s.shadow.inherit=False
        return s
    def CARD(sl,x,y,w,h,k="card_tint"):
        s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
        try: s.adjustments[0]=0.045
        except: pass
        s.fill.solid(); s.fill.fore_color.rgb=C(k)
        s.line.fill.background(); s.shadow.inherit=False
        return s
    def TXT(sl,text,x,y,w,h,sz,bold=False,k="text_dk",al=PP_ALIGN.LEFT,it=False):
        tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=al
        r=p.add_run(); r.text=str(text)
        r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=it
        r.font.name=FONT; r.font.color.rgb=C(k)
        return tb
    def CIRCLE(sl,cx,cy,d,label,fill_k="accent",text_k=None):
        s=sl.shapes.add_shape(MSO_SHAPE.OVAL,cx,cy,d,d)
        s.fill.solid(); s.fill.fore_color.rgb=C(fill_k)
        s.line.fill.background(); s.shadow.inherit=False
        tf=s.text_frame; tf.word_wrap=False
        tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=str(label)
        r.font.size=Pt(max(10,int(Emu(d).inches*40)))
        r.font.bold=True; r.font.name=FONT
        r.font.color.rgb=W() if not text_k else C(text_k)
        return s
    def DOTS(sl,items,x,y,w,h,sz=15,dot_k="accent",text_k="text_dk",gap=10):
        if not items: return
        tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
        for i,b in enumerate(items):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.space_before=Pt(gap); p.space_after=Pt(2)
            r1=p.add_run(); r1.text="●  "
            r1.font.size=Pt(sz-3); r1.font.name=FONT; r1.font.color.rgb=C(dot_k); r1.font.bold=True
            r2=p.add_run(); r2.text=str(b)
            r2.font.size=Pt(sz); r2.font.name=FONT; r2.font.color.rgb=C(text_k)
    def PAGENUM(sl,n,k="sub_dk"):
        tb=sl.shapes.add_textbox(W_-Inches(0.8),H_-Inches(0.45),Inches(0.6),Inches(0.3))
        p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
        r=p.add_run(); r.text=str(n)
        r.font.size=Pt(10); r.font.name=FONT; r.font.color.rgb=C(k)
    def CIRCLE_PLAIN(sl,x,y,d,k,blend_bg=None):
        s=sl.shapes.add_shape(MSO_SHAPE.OVAL,x,y,d,d)
        if blend_bg:
            base=pal[k]; mixed=tuple(int(base[i]*0.35+blend_bg[i]*0.65) for i in range(3))
            s.fill.solid(); s.fill.fore_color.rgb=RGBColor(*mixed)
        else:
            s.fill.solid(); s.fill.fore_color.rgb=C(k)
        s.line.fill.background(); s.shadow.inherit=False
        return s
    def DECOR_PANEL(sl,px,py,pw,ph,image_bytes=None):
        if image_bytes:
            sl.shapes.add_picture(io.BytesIO(image_bytes),px,py,pw,ph); return
        RECT(sl,px,py,pw,ph,"bg_dark2")
        cx=px+int(pw*0.62); cy=py+int(ph*0.42)
        CIRCLE_PLAIN(sl,cx-Inches(2.6),cy-Inches(2.6),Inches(5.2),"accent2",pal["bg_dark2"])
        CIRCLE_PLAIN(sl,cx-Inches(1.6),cy-Inches(1.6),Inches(3.2),"accent",pal["bg_dark2"])
        CIRCLE_PLAIN(sl,cx-Inches(0.5),cy-Inches(0.5),Inches(1.0),"accent2")

    def slice_bullets(buls):
        """숫자접두(01. 02.) 제거"""
        return [str(b).lstrip("0123456789. ").strip() for b in buls]

    for idx, si in enumerate(slides_data):
        sl = prs.slides.add_slide(BL)
        lay  = si.get("layout", "content")
        ttl  = si.get("title", "")
        sub  = si.get("subtitle", "")
        body = si.get("body", "")
        buls = si.get("bullets", []) or []
        hdrs = si.get("headers", []) or []
        rows = si.get("rows", []) or []
        sn   = idx + 1
        is_toc = lay == "content" and any(str(b).startswith(("01","02","1.","2.")) for b in buls)

        # ═══ DARK: title / closing / section ═══
        if lay in ("title", "closing", "section"):
            RECT(sl,0,0,W_,H_,"bg_dark")
            panel_w = int(W_*0.42)
            DECOR_PANEL(sl, W_-panel_w, 0, panel_w, H_,
                        cover_image_bytes if lay in ("title","closing") else None)
            if lay == "title":
                org=(summary or {}).get("주최기관","")
                if org: TXT(sl,org,Inches(0.6),Inches(0.6),Inches(6.5),Inches(0.5),14,True,k="accent2")
                TXT(sl,ttl,Inches(0.6),Inches(2.3),Inches(7.0),Inches(2.6),38,True,k="text_lt")
                if sub: TXT(sl,sub,Inches(0.6),Inches(5.1),Inches(7.0),Inches(1.0),14,k="sub_lt")
            elif lay == "closing":
                TXT(sl,ttl,Inches(0.6),Inches(2.6),Inches(7.0),Inches(1.6),42,True,k="text_lt")
                if sub: TXT(sl,sub,Inches(0.6),Inches(3.9),Inches(7.0),Inches(0.7),16,k="sub_lt")
                dept=" ".join(filter(None,[(summary or {}).get("담당부서",""),(summary or {}).get("담당자","")])).strip()
                if dept: TXT(sl,dept,Inches(0.6),Inches(6.7),Inches(7.0),Inches(0.4),11,k="sub_lt",it=True)
            elif lay == "section":
                CIRCLE(sl,Inches(0.6),Inches(2.6),Inches(1.1),f"{sn}","accent2")
                TXT(sl,ttl,Inches(0.6),Inches(4.0),Inches(7.5),Inches(1.6),32,True,k="text_lt")

        # ═══ TOC: Contents 분할형 (전용 레이아웃) ═══
        elif is_toc:
            RECT(sl,0,0,W_,H_,"bg_light")
            left_w = int(W_*0.38)
            RECT(sl,0,0,left_w,H_,"panel")
            tb = TXT(sl,"Contents",0,int(H_*0.44),left_w,Inches(0.9),36,True,k="text_lt",al=PP_ALIGN.CENTER)
            line_y = int(H_*0.44) + Inches(0.8)
            RECT(sl, int(left_w*0.30), line_y, int(left_w*0.40), Pt(1.2), "bg_light")

            items = slice_bullets(buls)[:6]
            right_x = left_w + Inches(0.5); right_w = W_ - left_w - Inches(0.9)
            row_h = min(Inches(0.98), (H_ - Inches(1.2)) // max(len(items),1))
            start_y = (H_ - row_h*len(items)) // 2
            for i, title in enumerate(items):
                ry = start_y + row_h*i
                TXT(sl,f"{i+1}",right_x,ry,Inches(0.9),Inches(0.8),36,True,k="num")
                TXT(sl,title,right_x+Inches(1.0),ry+Inches(0.14),right_w-Inches(1.0),Inches(0.55),19,True,k="text_dk")
                if i < len(items)-1:
                    RECT(sl,right_x,ry+row_h-Inches(0.06),right_w,Pt(1),"card_tint")
            PAGENUM(sl,sn)

        # ═══ LIGHT: content / two_column / highlight / table ═══
        else:
            RECT(sl,0,0,W_,H_,"bg_light")
            CIRCLE(sl,Inches(0.5),Inches(0.45),Inches(0.62),f"{sn}","accent")
            TXT(sl,ttl,Inches(1.35),Inches(0.42),Inches(10.5),Inches(0.6),26,True,k="text_dk")
            if body:
                TXT(sl,body,Inches(1.35),Inches(1.05),Inches(11.2),Inches(0.6),12.5,k="sub_dk")
                content_top = Inches(1.85)
            else:
                content_top = Inches(1.35)

            if lay == "highlight":
                stat_n=si.get("stat_number",""); stat_l=si.get("stat_label","")
                card_w=Inches(4.3); card_h=H_-content_top-Inches(0.5)
                CARD(sl,Inches(0.5),content_top,card_w,card_h)
                TXT(sl,stat_n,Inches(0.5),content_top+Inches(0.55),card_w,Inches(1.5),52,True,k="accent",al=PP_ALIGN.CENTER)
                TXT(sl,stat_l,Inches(0.5),content_top+Inches(2.0),card_w,Inches(0.5),14,k="sub_dk",al=PP_ALIGN.CENTER)
                DOTS(sl,buls,Inches(5.2),content_top+Inches(0.1),Inches(7.6),card_h,14)
                PAGENUM(sl,sn)

            elif lay == "two_column":
                lt=si.get("left_title","현황"); lb=si.get("left_bullets",[]) or []
                rt=si.get("right_title","목표"); rb=si.get("right_bullets",[]) or []
                cw=Inches(5.85); ch=H_-content_top-Inches(0.5)
                lx=Inches(0.5); rx=lx+cw+Inches(0.3)
                for x,ct,cb,bk in [(lx,lt,lb,"accent"),(rx,rt,rb,"accent2")]:
                    CARD(sl,x,content_top,cw,ch)
                    CIRCLE(sl,x+Inches(0.25),content_top+Inches(0.25),Inches(0.42),"",bk)
                    TXT(sl,ct,x+Inches(0.85),content_top+Inches(0.22),cw-Inches(1.1),Inches(0.5),16,True,k="text_dk")
                    DOTS(sl,cb,x+Inches(0.35),content_top+Inches(0.95),cw-Inches(0.6),ch-Inches(1.1),13,dot_k=bk)
                PAGENUM(sl,sn)

            elif lay=="table" and hdrs and rows:
                cc=len(hdrs); rc=len(rows)+1
                tbl=sl.shapes.add_table(rc,cc,Inches(0.5),content_top,Inches(12.3),H_-content_top-Inches(0.5)).table
                cw_=Inches(12.3)//cc
                for c in range(cc): tbl.columns[c].width=cw_
                for c,h in enumerate(hdrs):
                    cell=tbl.cell(0,c); cell.text=str(h)
                    cell.fill.solid(); cell.fill.fore_color.rgb=C("accent")
                    p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
                    r=p.runs[0] if p.runs else p.add_run()
                    r.font.bold=True; r.font.size=Pt(12); r.font.color.rgb=W(); r.font.name=FONT
                for ri,row in enumerate(rows):
                    rk="card" if ri%2==0 else "card_tint"
                    for c,v in enumerate(row[:cc]):
                        cell=tbl.cell(ri+1,c); cell.text=str(v)
                        cell.fill.solid(); cell.fill.fore_color.rgb=C(rk)
                        p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
                        r=p.runs[0] if p.runs else p.add_run()
                        r.font.size=Pt(11); r.font.color.rgb=C("text_dk"); r.font.name=FONT
                PAGENUM(sl,sn)

            else:  # 일반 content
                DOTS(sl,buls,Inches(0.6),content_top,Inches(11.8),H_-content_top-Inches(0.5),15)
                PAGENUM(sl,sn)

    buf=io.BytesIO(); prs.save(buf); return buf.getvalue()



def render_ppt():
    st.subheader("📊 PPT 자동생성")
    st.caption("행사 계획서 기반으로 현장에서 바로 발표 가능한 PPT를 생성합니다.")
    s = _summary_str()
    c1, c2 = st.columns([1, 2])
    n     = c1.slider("슬라이드 수", 8, 20, 12)
    theme = c2.selectbox("색상 테마", list(PALETTES.keys()))

    with st.expander("🖼️ 표지·클로징용 이미지 업로드 (선택)", expanded=False):
        st.caption("행사 포스터나 대표 사진을 올리면 표지/클로징 우측 절반에 실제 이미지가 들어갑니다. 없으면 자동 생성 그래픽이 대신 들어갑니다.")
        cover_img = st.file_uploader("이미지 (jpg/png)", type=["jpg", "jpeg", "png"], key="ppt_cover_img")
    cover_bytes = cover_img.read() if cover_img else None

    if st.button("🖥️ PPT 생성", type="primary"):
        with st.spinner("AI가 슬라이드 구성 중... (10~20초)"):
            raw = ai(_prompt_ppt(s, n))
        try:
            slides = json.loads(re.sub(r"```json|```", "", raw).strip())
        except Exception as e:
            st.error(f"파싱 오류: {e}")
            with st.expander("AI 원본 확인"):
                st.text(raw[:1000])
            return

        empty_cnt = sum(1 for si in slides if _is_empty(si))
        if empty_cnt:
            with st.spinner(f"⚠️ 빈 슬라이드 {empty_cnt}개 감지 → 자동 보완 중..."):
                slides = _validate_and_fix(slides, s)
            still_empty = sum(1 for si in slides if _is_empty(si))
            if still_empty:
                st.warning(f"⚠️ {still_empty}개 슬라이드는 내용을 직접 채워주세요.")

        with st.spinner("PPT 파일 생성 중..."):
            data = _build_pptx(slides, st.session_state.get("plan_summary_dict", {}), theme, cover_bytes)
        if data:
            name = st.session_state.get("plan_summary_dict", {}).get("행사명", "행사")
            st.session_state.update({"ppt_bytes": data,
                                     "ppt_name": f"{name}_발표자료.pptx",
                                     "ppt_count": len(slides)})

    if st.session_state.get("ppt_bytes"):
        st.success(f"✅ {st.session_state['ppt_count']}개 슬라이드 생성 완료!")
        st.download_button(
            "⬇️ PPT 다운로드 (.pptx)",
            data=st.session_state["ppt_bytes"],
            file_name=st.session_state["ppt_name"],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )
