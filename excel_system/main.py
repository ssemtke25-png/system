"""
tab7_hwpx_ai.py
행사 계획서 기반 AI 문서 자동생성 탭
- HWPX 업로드 → 요약 → 7가지 생성 기능
"""

import streamlit as st
import google.generativeai as genai
import zipfile
import os
import io
import re
import json
import base64
import subprocess
import tempfile
import requests
from pathlib import Path
from datetime import datetime

# ─────────────────────────────────────────────
# 초기 설정
# ─────────────────────────────────────────────
GEMINI_MODEL = "gemini-2.5-flash-lite"
MALGUN_GOTHIC = "맑은 고딕"


def get_gemini_model():
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    return genai.GenerativeModel(GEMINI_MODEL)


# ─────────────────────────────────────────────
# HWPX 텍스트 추출
# ─────────────────────────────────────────────
def extract_hwpx_text(hwpx_file) -> str:
    """HWPX 파일에서 텍스트 추출"""
    text_parts = []
    try:
        with zipfile.ZipFile(io.BytesIO(hwpx_file.read()), "r") as zf:
            xml_files = [n for n in zf.namelist() if n.endswith(".xml")]
            # Contents/ 폴더 우선
            content_xmls = sorted([f for f in xml_files if "Contents" in f or "content" in f.lower()])
            target = content_xmls if content_xmls else xml_files
            for name in target:
                try:
                    data = zf.read(name).decode("utf-8", errors="ignore")
                    # XML 태그 제거 후 텍스트 추출
                    clean = re.sub(r"<[^>]+>", " ", data)
                    clean = re.sub(r"\s+", " ", clean).strip()
                    if len(clean) > 50:
                        text_parts.append(clean)
                except Exception:
                    continue
    except Exception as e:
        return f"[파일 읽기 오류: {e}]"
    full = " ".join(text_parts)
    return full[:8000] if len(full) > 8000 else full  # 토큰 절약


def summarize_plan(raw_text: str) -> str:
    """계획서 원문 → 구조화 요약 (이후 모든 기능에 재사용)"""
    model = get_gemini_model()
    prompt = f"""다음은 행사 계획서 원문입니다. 아래 항목들을 JSON 형식으로 정리해주세요.
반드시 JSON만 출력하고 마크다운 코드블록 없이 순수 JSON만 반환하세요.

원문:
{raw_text}

출력 JSON 형식:
{{
  "행사명": "",
  "주최기관": "",
  "주관기관": "",
  "일시": "",
  "장소": "",
  "대상": "",
  "참석예정인원": "",
  "행사목적": "",
  "주요프로그램": [""],
  "예산": "",
  "담당부서": "",
  "담당자": "",
  "기타특이사항": ""
}}"""
    response = model.generate_content(prompt)
    return response.text.strip()


# ─────────────────────────────────────────────
# DOCX 생성 유틸
# ─────────────────────────────────────────────
def build_docx_via_node(js_code: str, out_name: str) -> bytes | None:
    """Node.js + docx 라이브러리로 DOCX 생성"""
    with tempfile.TemporaryDirectory() as td:
        js_path = Path(td) / "gen.mjs"
        out_path = Path(td) / out_name
        js_path.write_text(js_code, encoding="utf-8")
        # docx 패키지 설치
        subprocess.run(["npm", "install", "--prefix", td, "docx"], capture_output=True)
        result = subprocess.run(
            ["node", str(js_path)],
            capture_output=True, text=True, cwd=td
        )
        if result.returncode != 0:
            st.error(f"DOCX 생성 오류:\n{result.stderr[:500]}")
            return None
        if out_path.exists():
            return out_path.read_bytes()
    return None


def make_doc_js(title: str, content_paragraphs: list[dict], out_name: str, summary: dict) -> str:
    """범용 DOCX 생성 JS 코드 반환
    content_paragraphs: [{"text": ..., "heading": 1|2|None, "bold": bool}]
    """
    import json as _json
    paras_json = _json.dumps(content_paragraphs, ensure_ascii=False)
    return f"""
import {{ Document, Packer, Paragraph, TextRun, HeadingLevel, AlignmentType }} from '{{}}/node_modules/docx/build/index.js'.replace('{{}}', process.cwd());
import fs from 'fs';
import path from 'path';

// docx 모듈 경로 동적 로드
const docxPath = path.join(process.cwd(), 'node_modules', 'docx');
const {{ Document: Doc, Packer: Pkr, Paragraph: Para, TextRun: TR, HeadingLevel: HL, AlignmentType: AT }} = await import(docxPath);

const paragraphs = {paras_json};

const children = paragraphs.map(p => {{
  if (p.heading === 1) {{
    return new Para({{ heading: HL.HEADING_1, children: [new TR({{ text: p.text, font: "{MALGUN_GOTHIC}" }})] }});
  }} else if (p.heading === 2) {{
    return new Para({{ heading: HL.HEADING_2, children: [new TR({{ text: p.text, font: "{MALGUN_GOTHIC}" }})] }});
  }} else {{
    return new Para({{ children: [new TR({{ text: p.text, bold: !!p.bold, font: "{MALGUN_GOTHIC}", size: 24 }})] }});
  }}
}});

const doc = new Doc({{
  styles: {{
    default: {{ document: {{ run: {{ font: "{MALGUN_GOTHIC}", size: 24 }} }} }},
    paragraphStyles: [
      {{ id: "Heading1", name: "Heading 1", basedOn: "Normal", next: "Normal",
        run: {{ size: 36, bold: true, font: "{MALGUN_GOTHIC}", color: "1F3864" }},
        paragraph: {{ spacing: {{ before: 400, after: 200 }} }} }},
      {{ id: "Heading2", name: "Heading 2", basedOn: "Normal", next: "Normal",
        run: {{ size: 28, bold: true, font: "{MALGUN_GOTHIC}", color: "2E74B5" }},
        paragraph: {{ spacing: {{ before: 240, after: 120 }} }} }},
    ]
  }},
  sections: [{{
    properties: {{
      page: {{
        size: {{ width: 11906, height: 16838 }},
        margin: {{ top: 1440, right: 1080, bottom: 1440, left: 1080 }}
      }}
    }},
    children
  }}]
}});

const buffer = await Pkr.toBuffer(doc);
fs.writeFileSync('{out_name}', buffer);
console.log('done');
"""


# ─────────────────────────────────────────────
# 공통: AI → 텍스트 → DOCX 바이트
# ─────────────────────────────────────────────
def ai_to_docx(system_prompt: str, user_prompt: str, out_name: str) -> bytes | None:
    """AI 생성 텍스트를 DOCX로 변환"""
    model = get_gemini_model()
    full_prompt = system_prompt + "\n\n" + user_prompt
    response = model.generate_content(full_prompt)
    content = response.text.strip()

    # 텍스트 → 단락 파싱
    paragraphs = []
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            paragraphs.append({"text": line[2:], "heading": 1})
        elif line.startswith("## "):
            paragraphs.append({"text": line[3:], "heading": 2})
        elif line.startswith("**") and line.endswith("**"):
            paragraphs.append({"text": line[2:-2], "bold": True})
        else:
            paragraphs.append({"text": line})

    return _paragraphs_to_docx(paragraphs, out_name), content


def _paragraphs_to_docx(paragraphs: list, out_name: str) -> bytes | None:
    """단락 리스트 → DOCX bytes (Python docx 직접 사용)"""
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()
        # 문서 속성 설정 (보안 경고 감소)
        doc.core_properties.author = "AI 문서 자동생성 시스템"
        doc.core_properties.title = out_name.replace(".docx", "")
        doc.core_properties.language = "ko-KR"
        style = doc.styles["Normal"]
        style.font.name = MALGUN_GOTHIC
        style.font.size = Pt(11)

        # 제목 스타일
        for h_style, sz, color in [
            ("Heading 1", 18, (31, 56, 100)),
            ("Heading 2", 14, (46, 116, 181)),
        ]:
            s = doc.styles[h_style]
            s.font.name = MALGUN_GOTHIC
            s.font.size = Pt(sz)
            s.font.color.rgb = RGBColor(*color)

        # 여백 설정
        for section in doc.sections:
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(3.0)
            section.right_margin = Cm(2.5)

        for p in paragraphs:
            text = p.get("text", "")
            h = p.get("heading")
            bold = p.get("bold", False)

            if h == 1:
                para = doc.add_heading(text, level=1)
            elif h == 2:
                para = doc.add_heading(text, level=2)
            else:
                para = doc.add_paragraph()
                run = para.add_run(text)
                run.font.name = MALGUN_GOTHIC
                run.font.size = Pt(11)
                if bold:
                    run.bold = True

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        st.error("python-docx 패키지가 필요합니다: pip install python-docx")
        return None


# ─────────────────────────────────────────────
# 1. 문서 4종
# ─────────────────────────────────────────────
DOC_TYPES = {
    "요약보고서": {
        "prompt": "다음 행사 계획서 요약을 바탕으로 내부 결재용 요약보고서를 작성하세요. 마크다운 헤딩(#, ##)을 사용해 구조화하세요.",
        "filename": "요약보고서.docx",
    },
    "국장인사말": {
        "prompt": "다음 행사 계획서 요약을 바탕으로 국장급 인사말을 작성하세요. 격식 있고 품위 있는 문체로 400~600자 분량으로 작성하세요.",
        "filename": "국장인사말.docx",
    },
    "과장인사말": {
        "prompt": "다음 행사 계획서 요약을 바탕으로 과장급 인사말을 작성하세요. 친근하고 따뜻한 문체로 300~500자 분량으로 작성하세요.",
        "filename": "과장인사말.docx",
    },
    "보도자료": {
        "prompt": "다음 행사 계획서 요약을 바탕으로 언론 배포용 보도자료를 작성하세요. 마크다운 헤딩(#, ##)을 사용하고 육하원칙에 따라 작성하세요.",
        "filename": "보도자료.docx",
    },
}


def _render_one_doc(doc_name: str, cfg: dict, summary_str: str):
    """문서 1종 렌더링 - 텍스트 표시 + 복사용"""
    sess_key = f"doc4_text_{doc_name}"

    st.markdown(f"**{doc_name}**")

    if st.button(f"✍️ {doc_name} 생성", key=f"doc4_btn_{doc_name}"):
        with st.spinner(f"{doc_name} 생성 중..."):
            model = get_gemini_model()
            full_prompt = cfg["prompt"] + f"\n\n행사 계획서 요약:\n{summary_str}"
            response = model.generate_content(full_prompt)
            content = response.text.strip()
        st.session_state[sess_key] = content

    if st.session_state.get(sess_key):
        content = st.session_state[sess_key]
        with st.expander(f"📋 {doc_name} 미리보기 (전체 선택 후 복사)", expanded=True):
            st.text_area(
                label="",
                value=content,
                height=400,
                key=f"textarea_{doc_name}",
                help="전체 선택(Ctrl+A) 후 복사(Ctrl+C) 하세요",
            )


def render_doc4():
    st.subheader("📄 문서 4종 자동생성")
    st.caption("생성 후 텍스트를 복사해서 한글/워드에 붙여넣기 하세요.")
    summary = st.session_state.get("plan_summary_dict", {})
    summary_str = json.dumps(summary, ensure_ascii=False, indent=2)

    doc_list = list(DOC_TYPES.items())
    col1, col2 = st.columns(2)
    with col1:
        _render_one_doc(doc_list[0][0], doc_list[0][1], summary_str)
        st.markdown("---")
        _render_one_doc(doc_list[2][0], doc_list[2][1], summary_str)
    with col2:
        _render_one_doc(doc_list[1][0], doc_list[1][1], summary_str)
        st.markdown("---")
        _render_one_doc(doc_list[3][0], doc_list[3][1], summary_str)


# ─────────────────────────────────────────────
# 2. 사회자 멘트
# ─────────────────────────────────────────────
MC_TEMPLATE = [
    "개회선언",
    "국민의례",
    "내빈소개",
    "기관장인사말",
    "축사",
    "주요프로그램",
    "폐회선언",
]


def render_mc():
    st.subheader("🎤 사회자 멘트 자동생성")
    summary = st.session_state.get("plan_summary_dict", {})
    summary_str = json.dumps(summary, ensure_ascii=False, indent=2)

    st.info("고정 순서: " + " → ".join(MC_TEMPLATE))

    if st.button("✍️ 사회자 멘트 생성", key="mc_gen"):
        prompt = f"""다음 행사 계획서 요약과 순서에 맞는 사회자 멘트를 작성하세요.
각 순서마다 ## 헤딩으로 구분하고 실제 말할 멘트를 작성하세요.
순서: {', '.join(MC_TEMPLATE)}

행사 계획서 요약:
{summary_str}

각 순서마다 자연스럽고 격식 있는 사회자 멘트를 100~200자 내외로 작성하세요."""

        with st.spinner("사회자 멘트 생성 중..."):
            model = get_gemini_model()
            response = model.generate_content(prompt)
            content = response.text.strip()
        st.session_state["mc_content"] = content

    if st.session_state.get("mc_content"):
        with st.expander("📋 사회자 멘트 (전체 선택 후 복사)", expanded=True):
            st.text_area(
                label="",
                value=st.session_state["mc_content"],
                height=500,
                key="textarea_mc",
                help="전체 선택(Ctrl+A) 후 복사(Ctrl+C) 하세요",
            )


# ─────────────────────────────────────────────
# 3. 현수막 문구
# ─────────────────────────────────────────────
def render_banner():
    st.subheader("🪧 현수막 문구 시안")
    summary = st.session_state.get("plan_summary_dict", {})
    summary_str = json.dumps(summary, ensure_ascii=False, indent=2)

    n = st.slider("시안 수", 3, 5, 3)

    if st.button("✍️ 현수막 문구 생성", key="banner_gen"):
        prompt = f"""다음 행사 계획서 요약을 바탕으로 현수막 문구 시안을 {n}개 작성하세요.
각 시안은 짧고 임팩트 있게, 행사 목적과 분위기에 맞게 작성하세요.
형식: 번호. 문구 (20자 내외)

행사 계획서 요약:
{summary_str}"""

        with st.spinner("현수막 문구 생성 중..."):
            model = get_gemini_model()
            response = model.generate_content(prompt)
            content = response.text.strip()

        st.session_state["banner_content"] = content

    if st.session_state.get("banner_content"):
        with st.expander("📋 현수막 문구 (전체 선택 후 복사)", expanded=True):
            st.text_area(
                label="",
                value=st.session_state["banner_content"],
                height=300,
                key="textarea_banner",
                help="전체 선택(Ctrl+A) 후 복사(Ctrl+C) 하세요",
            )


# ─────────────────────────────────────────────
# 4. 결과보고서 초안
# ─────────────────────────────────────────────
def render_result_report():
    st.subheader("📋 결과보고서 초안")
    summary = st.session_state.get("plan_summary_dict", {})
    summary_str = json.dumps(summary, ensure_ascii=False, indent=2)

    col1, col2 = st.columns(2)
    with col1:
        actual_attendance = st.text_input("실제 참석 인원", placeholder="예: 150명")
    with col2:
        result_note = st.text_area("특이사항/결과 메모", placeholder="행사 결과 특이사항 입력", height=100)

    if st.button("✍️ 결과보고서 초안 생성", key="result_gen"):
        prompt = f"""다음 행사 계획서 요약을 바탕으로 행사 결과보고서 초안을 작성하세요.
마크다운 헤딩(#, ##)을 사용해 구조화하세요.
아래 항목을 반드시 포함하세요:
- 행사 개요
- 행사 결과 (실제 참석 인원: {actual_attendance or '미정'})
- 주요 내용 및 성과
- 특이사항: {result_note or '없음'}
- 향후 계획

행사 계획서 요약:
{summary_str}"""

        docx_bytes, content = ai_to_docx("행사 결과보고서를 작성하세요.", prompt, "결과보고서초안.docx")
        if content:
            st.session_state["result_content"] = content

    if st.session_state.get("result_content"):
        with st.expander("📋 결과보고서 초안 (전체 선택 후 복사)", expanded=True):
            st.text_area(
                label="",
                value=st.session_state["result_content"],
                height=500,
                key="textarea_result",
                help="전체 선택(Ctrl+A) 후 복사(Ctrl+C) 하세요",
            )


# ─────────────────────────────────────────────
# 5. PPT 자동생성
# ─────────────────────────────────────────────

PPT_PROMPT = """당신은 공공기관 행사 PPT를 전문으로 만드는 전문가입니다.
아래 행사 계획서 요약을 바탕으로 현장에서 바로 사용할 수 있는 발표자료 슬라이드를 구성하세요.

[슬라이드 구성 원칙]
1. 첫 슬라이드: 반드시 "title" 레이아웃 (행사명 + 일시/장소/주최 부제목)
2. 두 번째 슬라이드: "content" 레이아웃으로 목차(Contents) 구성
3. 중간 섹션 전환 시: "section" 레이아웃 사용
4. 마지막 슬라이드: "closing" 레이아웃 (감사합니다 + 기관 정보)
5. "table" 레이아웃: 일정표나 예산표처럼 표 형태가 필요할 때 사용

[필수 포함 슬라이드 - 이 순서로 구성]
- 표지 (title)
- 목차 (content)
- 행사 개요 (content): 목적, 일시, 장소, 대상, 주최/주관
- 주요 프로그램 (table): 시간표 형식으로
- 세부 내용 슬라이드들 (계획서 내용에 따라 구성)
- 마무리/기대효과 (content)
- 클로징 (closing)

[bullets 작성 규칙]
- 한 슬라이드당 bullets 최대 6개
- 각 bullet은 30자 이내로 간결하게
- 핵심 키워드 위주, 불필요한 조사/어미 생략
- table 레이아웃의 경우 headers와 rows를 별도 제공

[출력 형식 - 반드시 순수 JSON만, 마크다운 없이]
[
  {
    "slide_num": 1,
    "layout": "title",
    "title": "행사명",
    "subtitle": "2026. 2. | 영상회의 | 국토교통부"
  },
  {
    "slide_num": 2,
    "layout": "content",
    "title": "목  차",
    "bullets": ["01. 행사 개요", "02. 주요 프로그램", "03. 세부 내용", "04. 기대효과"]
  },
  {
    "slide_num": 3,
    "layout": "table",
    "title": "행사 일정",
    "headers": ["시간", "내용", "비고"],
    "rows": [["09:00~09:30", "등록 및 접수", "로비"], ["09:30~09:40", "개회선언", "사회자"]]
  },
  {
    "slide_num": 99,
    "layout": "closing",
    "title": "감사합니다",
    "subtitle": "주최기관명"
  }
]

슬라이드 수: {slide_count}개 (표지·클로징 포함)

행사 계획서 요약:
{summary_str}"""


def render_ppt():
    st.subheader("📊 PPT 자동생성")
    st.caption("행사 계획서 기반으로 현장 사용 가능한 발표자료를 생성합니다.")
    summary = st.session_state.get("plan_summary_dict", {})
    summary_str = json.dumps(summary, ensure_ascii=False, indent=2)

    col1, col2 = st.columns([1, 2])
    with col1:
        slide_count = st.slider("슬라이드 수", 8, 20, 12)
    with col2:
        theme = st.selectbox(
            "색상 테마",
            ["네이비 (공공기관 정장)", "다크블루 (격식)", "그린 (환경/생태)", "버건디 (품격)"],
        )

    if st.button("🖥️ PPT 생성", key="ppt_gen", type="primary"):
        prompt = PPT_PROMPT.format(slide_count=slide_count, summary_str=summary_str)

        with st.spinner("AI가 슬라이드 내용을 구성 중... (10~20초 소요)"):
            model = get_gemini_model()
            response = model.generate_content(prompt)
            raw = response.text.strip()

        try:
            clean = re.sub(r"```json|```", "", raw).strip()
            slides_data = json.loads(clean)
        except Exception as e:
            st.error(f"슬라이드 데이터 파싱 오류: {e}")
            with st.expander("AI 원본 응답 확인"):
                st.text(raw[:1000])
            return

        with st.spinner("PPT 파일 생성 중..."):
            pptx_bytes = _build_pptx(slides_data, summary, theme)

        if pptx_bytes:
            event_name = summary.get("행사명", "행사") if summary else "행사"
            st.session_state["ppt_bytes"] = pptx_bytes
            st.session_state["ppt_name"] = f"{event_name}_발표자료.pptx"
            st.session_state["ppt_count"] = len(slides_data)

    if st.session_state.get("ppt_bytes"):
        st.success(f"✅ {st.session_state.get('ppt_count', '')}개 슬라이드 생성 완료!")
        st.download_button(
            "⬇️ PPT 다운로드 (.pptx)",
            data=st.session_state["ppt_bytes"],
            file_name=st.session_state.get("ppt_name", "발표자료.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
        )


# 테마별 색상 팔레트
THEMES = {
    "네이비 (공공기관 정장)": {
        "dark":    (0x1F, 0x38, 0x64),   # 진남색
        "accent":  (0x2E, 0x74, 0xB5),   # 파랑
        "accent2": (0x70, 0xAD, 0x47),   # 초록 포인트
        "light_bg":(0xF0, 0xF4, 0xF9),
        "text":    (0x1A, 0x1A, 0x2E),
        "sub_text":(0xBD, 0xD7, 0xEE),
    },
    "다크블루 (격식)": {
        "dark":    (0x0D, 0x1B, 0x2A),
        "accent":  (0x1B, 0x4F, 0x72),
        "accent2": (0xF3, 0x9C, 0x12),
        "light_bg":(0xF2, 0xF3, 0xF4),
        "text":    (0x17, 0x20, 0x2A),
        "sub_text":(0xAB, 0xC4, 0xD8),
    },
    "그린 (환경/생태)": {
        "dark":    (0x1E, 0x4D, 0x2B),
        "accent":  (0x27, 0xAE, 0x60),
        "accent2": (0xF3, 0x9C, 0x12),
        "light_bg":(0xF0, 0xF9, 0xF1),
        "text":    (0x1A, 0x2E, 0x1C),
        "sub_text":(0xA9, 0xD9, 0xB4),
    },
    "버건디 (품격)": {
        "dark":    (0x6B, 0x21, 0x2C),
        "accent":  (0xC0, 0x39, 0x4B),
        "accent2": (0xE8, 0xB4, 0x6A),
        "light_bg":(0xFD, 0xF2, 0xF2),
        "text":    (0x2E, 0x11, 0x16),
        "sub_text":(0xF5, 0xC6, 0xCC),
    },
}


def _build_pptx(slides_data: list, summary: dict, theme: str = "네이비 (공공기관 정장)") -> bytes | None:
    """python-pptx로 PPT 생성 - 완성도 높은 버전"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree

        palette = THEMES.get(theme, THEMES["네이비 (공공기관 정장)"])

        def rgb(key):
            return RGBColor(*palette[key])

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        W = prs.slide_width
        H = prs.slide_height
        blank = prs.slide_layouts[6]

        def add_rect(slide, x, y, w, h, color_key, line=False):
            from pptx.util import Emu
            shape = slide.shapes.add_shape(1, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(color_key)
            if line:
                shape.line.color.rgb = rgb(color_key)
                shape.line.width = Pt(0.5)
            else:
                shape.line.fill.background()
            return shape

        def add_text(slide, text, x, y, w, h, size, bold=False,
                     color_key="text", align=PP_ALIGN.LEFT, wrap=True, italic=False):
            txb = slide.shapes.add_textbox(x, y, w, h)
            tf = txb.text_frame
            tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = rgb(color_key)
            run.font.name = MALGUN_GOTHIC
            return txb

        def add_page_number(slide, num):
            """슬라이드 번호"""
            txb = slide.shapes.add_textbox(
                Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3)
            )
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = str(num)
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            run.font.name = MALGUN_GOTHIC

        for idx, slide_info in enumerate(slides_data):
            slide = prs.slides.add_slide(blank)
            layout = slide_info.get("layout", "content")
            title_text = slide_info.get("title", "")
            subtitle = slide_info.get("subtitle", "")
            bullets = slide_info.get("bullets", [])
            headers = slide_info.get("headers", [])
            rows = slide_info.get("rows", [])
            slide_num = idx + 1

            # ── title 슬라이드 ────────────────────────
            if layout == "title":
                # 전체 배경
                add_rect(slide, 0, 0, W, H, "dark")
                # 좌측 세로 액센트 바
                add_rect(slide, 0, 0, Inches(0.18), H, "accent")
                # 하단 장식 바
                add_rect(slide, 0, int(H * 0.78), W, int(H * 0.035), "accent")
                # 우하단 밝은 사각형 장식
                add_rect(slide, int(W * 0.75), int(H * 0.78), int(W * 0.25), int(H * 0.035), "accent2")

                # 주최기관 (상단 작게)
                org = summary.get("주최기관", "") if summary else ""
                if org:
                    add_text(slide, org, Inches(0.5), Inches(0.4), Inches(12), Inches(0.5),
                             size=13, color_key="sub_text", align=PP_ALIGN.CENTER)

                # 메인 제목
                add_text(slide, title_text,
                         Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.5),
                         size=38, bold=True, color_key="sub_text",
                         align=PP_ALIGN.CENTER)

                # 구분선
                add_rect(slide, Inches(4), Inches(4.5), Inches(5.3), Inches(0.04), "accent")

                # 부제목 (일시·장소 등)
                if subtitle:
                    add_text(slide, subtitle,
                             Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.8),
                             size=16, color_key="sub_text", align=PP_ALIGN.CENTER)

            # ── closing 슬라이드 ──────────────────────
            elif layout == "closing":
                add_rect(slide, 0, 0, W, H, "dark")
                add_rect(slide, 0, 0, Inches(0.18), H, "accent")
                add_rect(slide, 0, int(H * 0.78), W, int(H * 0.035), "accent")
                add_rect(slide, int(W * 0.75), int(H * 0.78), int(W * 0.25), int(H * 0.035), "accent2")

                add_text(slide, title_text,
                         Inches(0.5), Inches(2.2), Inches(12.3), Inches(2),
                         size=44, bold=True, color_key="sub_text", align=PP_ALIGN.CENTER)

                add_rect(slide, Inches(4), Inches(4.4), Inches(5.3), Inches(0.04), "accent")

                if subtitle:
                    add_text(slide, subtitle,
                             Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.8),
                             size=18, color_key="sub_text", align=PP_ALIGN.CENTER)

            # ── section 슬라이드 ──────────────────────
            elif layout == "section":
                add_rect(slide, 0, 0, W, H, "dark")
                add_rect(slide, 0, 0, Inches(0.18), H, "accent2")
                # 중앙 액센트 가로줄
                add_rect(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.06), "accent2")

                # 섹션 번호
                num_text = f"{slide_num - 1:02d}"
                add_text(slide, num_text,
                         Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                         size=80, bold=True, color_key="accent",
                         align=PP_ALIGN.CENTER)

                add_text(slide, title_text,
                         Inches(0.5), Inches(3.7), Inches(12.3), Inches(1.5),
                         size=30, bold=True, color_key="sub_text",
                         align=PP_ALIGN.CENTER)

            # ── table 슬라이드 ────────────────────────
            elif layout == "table" and headers and rows:
                add_rect(slide, 0, 0, W, H, "light_bg")
                add_rect(slide, 0, 0, W, Inches(1.15), "dark")
                add_rect(slide, 0, Inches(1.08), W, Inches(0.07), "accent")

                add_text(slide, title_text,
                         Inches(0.4), Inches(0.18), Inches(12), Inches(0.75),
                         size=24, bold=True, color_key="sub_text")

                # 표 생성
                col_count = len(headers)
                row_count = len(rows) + 1  # 헤더 포함
                tbl_left = Inches(0.5)
                tbl_top = Inches(1.3)
                tbl_w = Inches(12.3)
                tbl_h = Inches(5.6)

                table = slide.shapes.add_table(
                    row_count, col_count, tbl_left, tbl_top, tbl_w, tbl_h
                ).table

                col_w = tbl_w // col_count
                for c in range(col_count):
                    table.columns[c].width = col_w

                # 헤더 행
                for c, hdr in enumerate(headers):
                    cell = table.cell(0, c)
                    cell.text = hdr
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb("accent")
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0] if p.runs else p.add_run()
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.name = MALGUN_GOTHIC

                # 데이터 행
                for r, row_data in enumerate(rows):
                    for c, val in enumerate(row_data[:col_count]):
                        cell = table.cell(r + 1, c)
                        cell.text = str(val)
                        # 짝수/홀수 행 색상
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = (
                            RGBColor(0xFF, 0xFF, 0xFF) if r % 2 == 0
                            else rgb("light_bg")
                        )
                        p = cell.text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.runs[0] if p.runs else p.add_run()
                        run.font.size = Pt(13)
                        run.font.color.rgb = rgb("text")
                        run.font.name = MALGUN_GOTHIC

                add_page_number(slide, slide_num)

            # ── content 슬라이드 (기본) ───────────────
            else:
                add_rect(slide, 0, 0, W, H, "light_bg")
                # 상단 헤더
                add_rect(slide, 0, 0, W, Inches(1.15), "dark")
                # 헤더 하단 액센트선
                add_rect(slide, 0, Inches(1.08), W, Inches(0.07), "accent")
                # 좌측 세로 포인트 바 (콘텐츠 영역)
                add_rect(slide, Inches(0.35), Inches(1.35), Inches(0.06), Inches(5.8), "accent2")

                add_text(slide, title_text,
                         Inches(0.4), Inches(0.18), Inches(12), Inches(0.75),
                         size=24, bold=True, color_key="sub_text")

                if bullets:
                    # 목차 슬라이드인지 감지 (bullets가 번호로 시작)
                    is_toc = any(b.startswith(("01", "02", "1.", "2.")) for b in bullets)

                    txb = slide.shapes.add_textbox(
                        Inches(0.6), Inches(1.4), Inches(12.3), Inches(5.7)
                    )
                    tf = txb.text_frame
                    tf.word_wrap = True

                    for i, bullet in enumerate(bullets):
                        p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p2.space_before = Pt(10 if is_toc else 6)
                        p2.space_after = Pt(4)
                        run2 = p2.add_run()
                        run2.text = ("  " + bullet) if is_toc else f"  ▪  {bullet}"
                        run2.font.size = Pt(20 if is_toc else 17)
                        run2.font.bold = is_toc
                        run2.font.color.rgb = rgb("text")
                        run2.font.name = MALGUN_GOTHIC

                add_page_number(slide, slide_num)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    except ImportError:
        st.error("python-pptx 패키지가 필요합니다: pip install python-pptx")
        return None
    except Exception as e:
        st.error(f"PPT 생성 오류: {e}")
        import traceback
        st.text(traceback.format_exc())
        return None


# ─────────────────────────────────────────────
# 6. 명찰 자동생성 (AI 호출 없음)
# ─────────────────────────────────────────────
def render_namecard():
    st.subheader("🪪 명찰 자동생성")
    summary = st.session_state.get("plan_summary_dict", {})
    event_name = summary.get("행사명", "행사") if summary else "행사"

    st.info("엑셀 파일에 **이름**, **소속** 컬럼이 있어야 합니다.")

    col1, col2 = st.columns(2)
    with col1:
        excel_file = st.file_uploader("명단 엑셀 업로드", type=["xlsx", "xls"], key="namecard_excel")
    with col2:
        schedule_text = st.text_area(
            "뒷면 일정표 내용",
            placeholder="09:00 등록\n10:00 개회식\n...",
            height=150,
        )

    custom_event = st.text_input("행사명 (비워두면 계획서에서 자동)", value="", placeholder=event_name)
    final_event = custom_event.strip() or event_name

    if excel_file and st.button("🪪 명찰 생성", key="namecard_gen"):
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(excel_file.read()), data_only=True)
            ws = wb.active

            # 첫 행을 헤더로 읽기
            headers = [str(cell.value).strip() if cell.value is not None else "" for cell in ws[1]]

            # 이름/소속 컬럼 인덱스 자동 감지
            name_idx = next(
                (i for i, h in enumerate(headers) if any(k in h for k in ["이름", "성명", "name", "Name"])),
                None,
            )
            org_idx = next(
                (i for i, h in enumerate(headers) if any(k in h for k in ["소속", "기관", "org", "Org", "dept", "Dept"])),
                None,
            )

            if name_idx is None:
                st.error(f"'이름' 또는 '성명' 컬럼을 찾을 수 없습니다. 현재 컬럼: {headers}")
                return
            if org_idx is None:
                st.warning("'소속' 컬럼을 찾을 수 없습니다. 소속 없이 생성합니다.")

            persons = []
            for row in ws.iter_rows(min_row=2, values_only=True):
                name = str(row[name_idx]).strip() if name_idx is not None and row[name_idx] is not None else ""
                org = str(row[org_idx]).strip() if org_idx is not None and row[org_idx] is not None else ""
                if name and name.lower() != "none":
                    persons.append({"name": name, "org": org})

            if not persons:
                st.error("명단에 유효한 데이터가 없습니다.")
                return

            st.success(f"총 {len(persons)}명 처리 중...")
            docx_bytes = _build_namecard_docx(persons, final_event, schedule_text)
            if docx_bytes:
                st.download_button(
                    f"⬇️ 명찰 {len(persons)}명 다운로드",
                    data=docx_bytes,
                    file_name="명찰.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )

        except Exception as e:
            st.error(f"명단 처리 오류: {e}")


def _build_namecard_docx(persons: list, event_name: str, schedule_text: str) -> bytes | None:
    """명찰 DOCX 생성 (A4 세로, 1장에 2세트 앞+뒤)"""
    try:
        from docx import Document
        from docx.shared import Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        import copy

        doc = Document()

        for section in doc.sections:
            section.page_width = Cm(21)
            section.page_height = Cm(29.7)
            section.top_margin = Cm(1.5)
            section.bottom_margin = Cm(1.5)
            section.left_margin = Cm(2.0)
            section.right_margin = Cm(2.0)

        def add_card_front(name: str, org: str, is_first_on_page: bool):
            """앞면 카드 추가"""
            if not is_first_on_page:
                # 구분선
                hr = doc.add_paragraph()
                hr.paragraph_format.space_before = Pt(0)
                hr.paragraph_format.space_after = Pt(0)
                pPr = hr._p.get_or_add_pPr()
                pBdr = OxmlElement("w:pBdr")
                bottom = OxmlElement("w:bottom")
                bottom.set(qn("w:val"), "dashed")
                bottom.set(qn("w:sz"), "6")
                bottom.set(qn("w:space"), "1")
                bottom.set(qn("w:color"), "888888")
                pBdr.append(bottom)
                pPr.append(pBdr)

            # 행사명
            p_event = doc.add_paragraph()
            p_event.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_event.paragraph_format.space_before = Pt(24)
            p_event.paragraph_format.space_after = Pt(8)
            r = p_event.add_run(event_name)
            r.font.name = MALGUN_GOTHIC
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

            # 소속
            if org:
                p_org = doc.add_paragraph()
                p_org.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_org.paragraph_format.space_before = Pt(4)
                p_org.paragraph_format.space_after = Pt(4)
                r_org = p_org.add_run(org)
                r_org.font.name = MALGUN_GOTHIC
                r_org.font.size = Pt(18)
                r_org.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)

            # 이름
            p_name = doc.add_paragraph()
            p_name.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_name.paragraph_format.space_before = Pt(8)
            p_name.paragraph_format.space_after = Pt(24)
            r_name = p_name.add_run(name)
            r_name.font.name = MALGUN_GOTHIC
            r_name.font.size = Pt(36)
            r_name.font.bold = True
            r_name.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        def add_card_back():
            """뒷면 카드 추가 (일정표)"""
            p_title = doc.add_paragraph()
            p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_title.paragraph_format.space_before = Pt(24)
            p_title.paragraph_format.space_after = Pt(8)
            r = p_title.add_run("행사 일정")
            r.font.name = MALGUN_GOTHIC
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

            schedule_lines = schedule_text.strip().split("\n") if schedule_text.strip() else ["일정 미정"]
            for line in schedule_lines:
                p_s = doc.add_paragraph()
                p_s.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_s.paragraph_format.space_before = Pt(2)
                p_s.paragraph_format.space_after = Pt(2)
                r_s = p_s.add_run(line.strip())
                r_s.font.name = MALGUN_GOTHIC
                r_s.font.size = Pt(14)

            # 마무리 여백
            p_end = doc.add_paragraph()
            p_end.paragraph_format.space_after = Pt(24)

        # 2명씩 묶어서 처리 (1장에 앞면2+뒷면2)
        for i in range(0, len(persons), 2):
            batch = persons[i: i + 2]

            # 앞면 (1~2번)
            add_card_front(batch[0]["name"], batch[0]["org"], True)
            if len(batch) > 1:
                add_card_front(batch[1]["name"], batch[1]["org"], False)

            # 페이지 나누기
            doc.add_page_break()

            # 뒷면 (2개)
            add_card_back()
            if len(batch) > 1:
                # 구분선 추가 후 두 번째 뒷면
                hr = doc.add_paragraph()
                hr.paragraph_format.space_before = Pt(0)
                hr.paragraph_format.space_after = Pt(0)
                pPr = hr._p.get_or_add_pPr()
                pBdr = OxmlElement("w:pBdr")
                bottom = OxmlElement("w:bottom")
                bottom.set(qn("w:val"), "dashed")
                bottom.set(qn("w:sz"), "6")
                bottom.set(qn("w:space"), "1")
                bottom.set(qn("w:color"), "888888")
                pBdr.append(bottom)
                pPr.append(pBdr)
                add_card_back()

            # 다음 배치 전 페이지 나누기 (마지막 제외)
            if i + 2 < len(persons):
                doc.add_page_break()

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    except ImportError:
        st.error("python-docx 패키지가 필요합니다: pip install python-docx")
        return None
    except Exception as e:
        st.error(f"명찰 생성 오류: {e}")
        return None


# ─────────────────────────────────────────────
# 7. 행사장 약도 (카카오맵 Static API)
# ─────────────────────────────────────────────
def render_map():
    st.subheader("🗺️ 행사장 약도")
    summary = st.session_state.get("plan_summary_dict", {})
    default_place = summary.get("장소", "") if summary else ""

    place_name = st.text_input("장소명 입력", value=default_place, placeholder="예: 서울시청 다목적홀")

    col1, col2, col3 = st.columns(3)
    with col1:
        zoom = st.slider("지도 확대 수준", 1, 14, 4)
    with col2:
        map_w = st.number_input("가로 (px)", value=800, min_value=200, max_value=1200)
    with col3:
        map_h = st.number_input("세로 (px)", value=600, min_value=200, max_value=900)

    if st.button("🗺️ 약도 생성", key="map_gen") and place_name:
        try:
            kakao_key = st.secrets["KAKAO_API_KEY"]

            # 1. 주소 → 좌표 변환
            geo_url = "https://dapi.kakao.com/v2/local/search/keyword.json"
            headers = {"Authorization": f"KakaoAK {kakao_key}"}
            geo_resp = requests.get(geo_url, headers=headers, params={"query": place_name})
            geo_data = geo_resp.json()

            if not geo_data.get("documents"):
                st.error(f"'{place_name}' 위치를 찾을 수 없습니다.")
                return

            doc0 = geo_data["documents"][0]
            lng = float(doc0["x"])
            lat = float(doc0["y"])
            found_name = doc0.get("place_name", place_name)

            st.info(f"📍 찾은 장소: **{found_name}** ({lat:.4f}, {lng:.4f})")

            # 2. Static Map 이미지 요청
            map_url = "https://map.kakao.com/v1/map/staticmap.png"
            params = {
                "center": f"{lng},{lat}",
                "level": zoom,
                "width": int(map_w),
                "height": int(map_h),
                "marker": f"DYNAMICMAP,2,{lng},{lat}",
            }
            map_resp = requests.get(map_url, headers=headers, params=params)

            if map_resp.status_code == 200 and "image" in map_resp.headers.get("Content-Type", ""):
                img_bytes = map_resp.content
                st.image(img_bytes, caption=f"📍 {found_name}", use_container_width=True)
                st.download_button(
                    "⬇️ 약도 이미지 다운로드",
                    data=img_bytes,
                    file_name=f"행사장약도_{found_name}.png",
                    mime="image/png",
                )
            else:
                # 대체: 카카오맵 링크 제공
                kakao_map_link = f"https://map.kakao.com/?q={requests.utils.quote(place_name)}"
                st.warning("Static Map API 응답 오류. 카카오맵 링크를 제공합니다.")
                st.markdown(f"[🗺️ 카카오맵에서 보기]({kakao_map_link})")
                st.write(f"응답 코드: {map_resp.status_code}")

        except KeyError:
            st.error("KAKAO_API_KEY가 st.secrets에 설정되어 있지 않습니다.")
        except Exception as e:
            st.error(f"지도 생성 오류: {e}")


# ─────────────────────────────────────────────
# 메인 탭 렌더러
# ─────────────────────────────────────────────
def render_tab7():
    st.title("🤖 AI 문서 자동생성")
    st.caption("HWPX 계획서를 업로드하면 각종 문서를 자동으로 생성합니다.")

    # ── 계획서 업로드 & 요약 ──────────────────────
    st.markdown("---")
    st.subheader("📁 행사 계획서 업로드")

    hwpx_file = st.file_uploader(
        "HWPX 파일 업로드 (.hwpx)",
        type=["hwpx", "hwp"],
        key="hwpx_upload",
    )

    if hwpx_file:
        # 파일명+크기로 동일 파일 여부 판단 (업로더가 매번 새 객체 반환하는 문제 방지)
        file_key = f"{hwpx_file.name}_{hwpx_file.size}"
        already_summarized = (
            st.session_state.get("hwpx_file_key") == file_key
            and "plan_summary_raw" in st.session_state
        )

        if not already_summarized:
            st.session_state["hwpx_file_key"] = file_key
            with st.spinner("계획서 텍스트 추출 및 요약 중..."):
                raw_text = extract_hwpx_text(hwpx_file)
                if len(raw_text) < 100:
                    st.error("계획서에서 텍스트를 충분히 추출하지 못했습니다.")
                else:
                    summary_raw = summarize_plan(raw_text)
                    st.session_state["plan_summary_raw"] = summary_raw
                    try:
                        clean = re.sub(r"```json|```", "", summary_raw).strip()
                        st.session_state["plan_summary_dict"] = json.loads(clean)
                    except Exception:
                        st.session_state["plan_summary_dict"] = {}

        if "plan_summary_dict" in st.session_state:
            summary = st.session_state["plan_summary_dict"]
            with st.expander("📋 계획서 요약 확인 (모든 기능에 재사용됨)", expanded=True):
                if summary:
                    cols = st.columns(2)
                    keys = list(summary.keys())
                    half = len(keys) // 2
                    for i, k in enumerate(keys):
                        v = summary[k]
                        if isinstance(v, list):
                            v = ", ".join(v)
                        with cols[0 if i < half else 1]:
                            st.markdown(f"**{k}:** {v}")
                else:
                    st.text(st.session_state.get("plan_summary_raw", "요약 없음"))
    else:
        st.info("HWPX 계획서를 업로드하면 아래 기능들이 활성화됩니다.")
        # 빈 딕셔너리로 기능은 계속 사용 가능하게
        if "plan_summary_dict" not in st.session_state:
            st.session_state["plan_summary_dict"] = {}

    # ── 기능 탭 ──────────────────────────────────
    st.markdown("---")
    tabs = st.tabs([
        "📄 문서 4종",
        "🎤 사회자 멘트",
        "🪧 현수막",
        "📋 결과보고서",
        "📊 PPT",
        "🪪 명찰",
        "🗺️ 행사장 약도",
    ])

    with tabs[0]:
        render_doc4()
    with tabs[1]:
        render_mc()
    with tabs[2]:
        render_banner()
    with tabs[3]:
        render_result_report()
    with tabs[4]:
        render_ppt()
    with tabs[5]:
        render_namecard()
    with tabs[6]:
        render_map()


# ─────────────────────────────────────────────
# 직접 실행 시
# ─────────────────────────────────────────────
if __name__ == "__main__":
    st.set_page_config(page_title="AI 문서 자동생성", page_icon="🤖", layout="wide")
    render_tab7()
